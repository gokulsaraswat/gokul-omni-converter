from __future__ import annotations

import json
import os
import shutil
import tempfile
import threading
import zipfile
from functools import partial
from http.server import ThreadingHTTPServer, SimpleHTTPRequestHandler
from datetime import datetime, timedelta
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from reportlab.pdfgen import canvas

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - local compatibility fallback
    from PyPDF2 import PdfReader  # type: ignore

from app_state import APP_NAME, APP_STATE_PATH, AppStateStore
from recovery_support import latest_state_backup, state_backup_dir
from about_profile import load_about_profile
from automation_core import (
    add_fingerprints,
    bundle_paths_as_zip,
    discover_watch_candidates,
    export_presets_to_json,
    import_presets_from_json,
    write_run_report,
)
from build_support import (
    export_diagnostics_report,
    export_state_snapshot,
    import_state_snapshot,
    export_text_file,
    render_activity_report_html,
    export_support_bundle,
)
from release_support import build_example_update_manifest, check_for_updates, export_workspace_bundle, import_workspace_bundle
from asset_support import (
    asset_cache_root,
    cache_summary as asset_cache_summary,
    clear_asset_cache,
    download_text_file,
    load_asset_config,
    resolve_local_or_remote_asset,
    save_asset_config,
)
from engagement_core import ensure_install_date, should_show_first_launch_splash, should_show_login_popup
from converter_core import (
    BatchConfig,
    PdfToolConfig,
    ENGINE_PURE_PYTHON,
    MODE_ANY_TO_PDF,
    MODE_DOCS_TO_PDF,
    MODE_HTML_TO_DOCX,
    MODE_HTML_TO_MD,
    MODE_HTML_TO_PDF,
    MODE_MD_TO_DOCX,
    MODE_MD_TO_HTML,
    MODE_MD_TO_PDF,
    MODE_PDF_TO_DOCX,
    MODE_PDF_TO_HTML,
    MODE_PDF_TO_IMAGES,
    MODE_PDF_TO_PPTX,
    MODE_PDF_TO_XLSX,
    MODE_PRESENTATIONS_TO_IMAGES,
    MODE_PRESENTATIONS_TO_PDF,
    MODE_SHEETS_TO_PDF,
    MODE_TEXT_TO_PDF,
    PDF_TOOL_COMPRESS,
    PDF_TOOL_EDIT_METADATA,
    PDF_TOOL_EXTRACT_PAGES,
    PDF_TOOL_IMAGE_OVERLAY,
    PDF_TOOL_LOCK,
    PDF_TOOL_MERGE,
    PDF_TOOL_REDACT_TEXT,
    PDF_TOOL_REDACT_AREA,
    PDF_TOOL_EDIT_TEXT,
    PDF_TOOL_REMOVE_PAGES,
    PDF_TOOL_REORDER_PAGES,
    PDF_TOOL_SIGN_VISIBLE,
    PDF_TOOL_SPLIT_EVERY_N,
    PDF_TOOL_SPLIT_RANGES,
    PDF_TOOL_TEXT_OVERLAY,
    PDF_TOOL_UNLOCK,
    PDF_TOOL_WATERMARK_IMAGE,
    PDF_TOOL_WATERMARK_TEXT,
    dependency_status,
    build_conversion_route_preview,
    process_batch,
    process_pdf_tool,
)


from mail_core import SMTPSettings, build_email_message, build_eml_draft, create_mailto_url
from ocr_core import OcrConfig, OcrError, detect_tesseract_status, extract_text_with_ocr, image_to_searchable_pdf, pdf_to_searchable_pdf
from organizer_core import (
    build_default_sequence,
    duplicate_positions,
    export_pages_as_images as organizer_export_pages_as_images,
    extract_selected_pdf as organizer_extract_selected_pdf,
    move_positions_down,
    move_positions_to_index,
    pdf_summary as organizer_pdf_summary,
    rotate_positions,
    save_sequence_as_pdf as organizer_save_sequence_as_pdf,
    sequence_from_payload,
    sequence_to_payload,
)


from link_ingest import cache_root_from_setting, download_many_urls, extract_urls
from workflow_support import directory_stats, prune_directory
from preview_support import render_preview
from responsive_ui import resolve_flow_layout_width
from ui_text import format_engine_label, format_flag, humanize_identifier


def validate_ui_text_helpers() -> None:
    if humanize_identifier("pure_python") != "Pure Python":
        raise AssertionError("pure_python should be humanized.")
    if humanize_identifier("top-left") != "Top Left":
        raise AssertionError("top-left should be humanized.")
    if format_engine_label("libre_office") != "Libre Office":
        raise AssertionError("Engine labels should be humanized.")
    if format_flag(True) != "Yes" or format_flag(False) != "No":
        raise AssertionError("Boolean flags should map to Yes/No.")


def validate_flow_wrap_helper() -> None:
    width = resolve_flow_layout_width(320, 640, 900, min_width=180)
    if width != 320:
        raise AssertionError("Flow width should prefer the real widget width.")
    width = resolve_flow_layout_width(0, 260, 900, min_width=180)
    if width != 260:
        raise AssertionError("Flow width should fall back to parent width when needed.")
    width = resolve_flow_layout_width(0, 0, 150, min_width=180)
    if width != 180:
        raise AssertionError("Flow width should honor the minimum width.")


def validate_preview_support(sample: dict[str, Path]) -> None:
    preview_targets = [
        sample["pdf"],
        sample["pdf_multi"],
        sample["img1"],
        sample["txt"],
        sample["md"],
        sample["html"],
        sample["docx"],
        sample["xlsx"],
    ]
    if "pptx" in sample:
        preview_targets.append(sample["pptx"])

    for path in preview_targets:
        result = render_preview(path)
        if result.image.width <= 0 or result.image.height <= 0:
            raise AssertionError(f"Preview rendering returned an empty image for {path.name}")
        if not result.summary.strip():
            raise AssertionError(f"Preview rendering returned an empty summary for {path.name}")

    pdf_page = render_preview(sample["pdf_multi"], page=3, zoom=1.25)
    if pdf_page.page_count < 4 or pdf_page.current_page != 3:
        raise AssertionError("PDF preview did not preserve page navigation metadata.")

    missing_path = sample["pdf"].with_name("missing_preview_target.pdf")
    missing_result = render_preview(missing_path)
    if "missing" not in missing_result.kind.lower() and "missing" not in missing_result.summary.lower():
        raise AssertionError("Missing-file preview fallback did not trigger as expected.")


def create_sample_files(root: Path) -> dict[str, Path]:
    paths: dict[str, Path] = {}

    img1 = root / "image_1.png"
    img2 = root / "image_2.png"
    for index, path in enumerate([img1, img2], start=1):
        image = Image.new("RGB", (900, 600), "white")
        draw = ImageDraw.Draw(image)
        draw.text((40, 40), f"Sample Image {index}", fill="black")
        image.save(path)
    paths["img1"] = img1
    paths["img2"] = img2

    ocr_image = root / "ocr_sample.png"
    ocr_canvas = Image.new("RGB", (1500, 420), "white")
    ocr_draw = ImageDraw.Draw(ocr_canvas)
    try:
        ocr_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 86)
    except Exception:
        ocr_font = ImageFont.load_default()
    ocr_draw.text((70, 120), "INVOICE 12345", fill="black", font=ocr_font)
    ocr_canvas.save(ocr_image)
    paths["ocr_image"] = ocr_image

    ocr_pdf = root / "ocr_sample.pdf"
    ocr_doc = fitz.open()
    ocr_page = ocr_doc.new_page(width=1500, height=420)
    ocr_page.insert_image(ocr_page.rect, filename=str(ocr_image))
    ocr_doc.save(str(ocr_pdf))
    ocr_doc.close()
    paths["ocr_pdf"] = ocr_pdf

    docx_path = root / "sample.docx"
    doc = Document()
    doc.add_heading("Sample DOCX", level=0)
    doc.add_paragraph("Hello from the sample DOCX file.")
    doc.add_paragraph("First bullet item", style="List Bullet")
    doc.add_paragraph("Second bullet item", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value 1"
    table.cell(1, 1).text = "Value 2"
    doc.save(docx_path)
    paths["docx"] = docx_path

    xlsx_path = root / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 10])
    ws.append(["Beta", 20])
    summary = wb.create_sheet("Summary")
    summary.append(["Metric", "Result"])
    summary.append(["Rows", 2])
    summary.append(["Status", "Ready"])
    wb.save(xlsx_path)
    paths["xlsx"] = xlsx_path

    md_path = root / "sample.md"
    md_path.write_text("# Sample Markdown\n\n- one\n- two\n\n**bold text**\n\n> quoted line\n\n```python\nprint(\'hi\')\n```\n", encoding="utf-8")
    paths["md"] = md_path

    txt_path = root / "sample.txt"
    txt_path.write_text("Plain text file\nSecond line\nThird line\n", encoding="utf-8")
    paths["txt"] = txt_path

    html_path = root / "sample.html"
    html_path.write_text(
        """<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8" /><title>Sample HTML</title></head>
<body>
  <h1>Sample HTML</h1>
  <p>This is a <strong>sample</strong> HTML document for conversion.</p>
  <ul><li>Alpha</li><li>Beta</li></ul>
  <table>
    <tr><th>Name</th><th>Value</th></tr>
    <tr><td>Row A</td><td>10</td></tr>
    <tr><td>Row B</td><td>20</td></tr>
  </table>
</body>
</html>
""",
        encoding="utf-8",
    )
    paths["html"] = html_path

    pdf_path = root / "sample.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 780, "Sample PDF for extraction")
    pdf.drawString(72, 760, "Row 1    Value 10")
    pdf.drawString(72, 740, "Row 2    Value 20")
    pdf.save()
    paths["pdf"] = pdf_path

    multi_pdf_path = root / "sample_multi.pdf"
    multi = canvas.Canvas(str(multi_pdf_path))
    for page_number in range(1, 7):
        multi.setFont("Helvetica", 18)
        multi.drawString(72, 780, f"Sample multi-page PDF page {page_number}")
        multi.setFont("Helvetica", 12)
        multi.drawString(72, 748, f"Page marker: {page_number}")
        multi.drawString(72, 728, f"Data row {page_number}A    Value {page_number * 10}")
        multi.drawString(72, 708, f"Data row {page_number}B    Value {page_number * 20}")
        multi.showPage()
    multi.save()
    paths["pdf_multi"] = multi_pdf_path

    watermark_path = root / "watermark.png"
    watermark = Image.new("RGBA", (700, 220), (255, 255, 255, 0))
    draw = ImageDraw.Draw(watermark)
    draw.rounded_rectangle((10, 10, 690, 210), radius=24, fill=(30, 144, 255, 100))
    draw.text((80, 80), "GOKUL OMNI", fill=(255, 255, 255, 210))
    watermark.save(watermark_path)
    paths["watermark"] = watermark_path

    try:
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = root / "sample_presentation.pptx"
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        textbox.text_frame.text = "Sample PPTX Slide"
        slide.shapes.add_picture(str(img1), Inches(1), Inches(2), width=Inches(5.5))
        presentation.save(str(pptx_path))
        paths["pptx"] = pptx_path
    except Exception:
        pass

    return paths


def build_conversion_jobs(sample: dict[str, Path]) -> tuple[list[BatchConfig], list[str]]:
    deps = dependency_status()
    jobs: list[BatchConfig] = []
    skipped: list[str] = []

    jobs.append(
        BatchConfig(
            mode=MODE_ANY_TO_PDF,
            files=[sample["img1"], sample["img2"], sample["txt"]],
            output_dir=Path("outputs") / "any_to_pdf",
            merge_to_one_pdf=True,
            merged_output_name="combined",
        )
    )
    jobs.append(BatchConfig(mode=MODE_TEXT_TO_PDF, files=[sample["txt"]], output_dir=Path("outputs") / "text_to_pdf"))
    jobs.append(
        BatchConfig(
            mode=MODE_PDF_TO_IMAGES,
            files=[sample["pdf"]],
            output_dir=Path("outputs") / "pdf_to_images",
            image_format="png",
            image_scale=1.5,
        )
    )
    jobs.append(BatchConfig(mode=MODE_PDF_TO_DOCX, files=[sample["pdf"]], output_dir=Path("outputs") / "pdf_to_docx"))
    jobs.append(BatchConfig(mode=MODE_PDF_TO_XLSX, files=[sample["pdf"]], output_dir=Path("outputs") / "pdf_to_xlsx"))
    jobs.append(BatchConfig(mode=MODE_PDF_TO_HTML, files=[sample["pdf_multi"]], output_dir=Path("outputs") / "pdf_to_html"))
    jobs.append(
        BatchConfig(
            mode=MODE_PDF_TO_PPTX,
            files=[sample["pdf"]],
            output_dir=Path("outputs") / "pdf_to_pptx",
            image_scale=1.7,
        )
    )
    jobs.append(BatchConfig(mode=MODE_HTML_TO_DOCX, files=[sample["html"]], output_dir=Path("outputs") / "html_to_docx"))
    jobs.append(BatchConfig(mode=MODE_HTML_TO_MD, files=[sample["html"]], output_dir=Path("outputs") / "html_to_md"))
    jobs.append(BatchConfig(mode=MODE_MD_TO_PDF, files=[sample["md"]], output_dir=Path("outputs") / "md_to_pdf_explicit"))
    jobs.append(BatchConfig(mode=MODE_MD_TO_HTML, files=[sample["md"]], output_dir=Path("outputs") / "md_to_html"))
    jobs.append(BatchConfig(mode=MODE_MD_TO_DOCX, files=[sample["md"]], output_dir=Path("outputs") / "md_to_docx"))
    jobs.append(BatchConfig(mode=MODE_HTML_TO_PDF, files=[sample["html"]], output_dir=Path("outputs") / "html_to_pdf_explicit"))

    # Pure-Python engine checks added in Patch 5.
    jobs.append(
        BatchConfig(
            mode=MODE_DOCS_TO_PDF,
            files=[sample["docx"]],
            output_dir=Path("outputs") / "docx_to_pdf_pure_python",
            engine_mode=ENGINE_PURE_PYTHON,
        )
    )
    jobs.append(
        BatchConfig(
            mode=MODE_SHEETS_TO_PDF,
            files=[sample["xlsx"]],
            output_dir=Path("outputs") / "sheets_to_pdf_pure_python",
            engine_mode=ENGINE_PURE_PYTHON,
        )
    )
    jobs.append(
        BatchConfig(
            mode=MODE_TEXT_TO_PDF,
            files=[sample["md"]],
            output_dir=Path("outputs") / "markdown_to_pdf_pure_python",
            engine_mode=ENGINE_PURE_PYTHON,
        )
    )
    jobs.append(
        BatchConfig(
            mode=MODE_TEXT_TO_PDF,
            files=[sample["html"]],
            output_dir=Path("outputs") / "html_to_pdf_pure_python",
            engine_mode=ENGINE_PURE_PYTHON,
        )
    )
    if sample.get("pptx"):
        jobs.append(
            BatchConfig(
                mode=MODE_PRESENTATIONS_TO_PDF,
                files=[sample["pptx"]],
                output_dir=Path("outputs") / "presentations_to_pdf_pure_python",
                engine_mode=ENGINE_PURE_PYTHON,
            )
        )
        jobs.append(
            BatchConfig(
                mode=MODE_PRESENTATIONS_TO_IMAGES,
                files=[sample["pptx"]],
                output_dir=Path("outputs") / "presentations_to_images",
                engine_mode=ENGINE_PURE_PYTHON,
                image_format="png",
                image_scale=1.4,
            )
        )

    if deps.get("LibreOffice"):
        jobs.append(
            BatchConfig(
                mode=MODE_DOCS_TO_PDF,
                files=[sample["docx"]],
                output_dir=Path("outputs") / "docx_to_pdf_libreoffice",
            )
        )
    else:
        skipped.append("LibreOffice not found. Optional high-fidelity Office conversion path was not smoke tested.")

    return jobs, skipped


def build_pdf_tool_jobs(sample: dict[str, Path]) -> list[PdfToolConfig]:
    return [
        PdfToolConfig(
            tool=PDF_TOOL_MERGE,
            files=[sample["pdf"], sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_merge",
            output_name="merged_tool_output",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_SPLIT_RANGES,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_split_ranges",
            page_spec="1-2; 3-4; 5-6",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_SPLIT_EVERY_N,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_split_every_n",
            every_n_pages=2,
        ),
        PdfToolConfig(
            tool=PDF_TOOL_EXTRACT_PAGES,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_extract",
            page_spec="1,3,5",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_REMOVE_PAGES,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_remove",
            page_spec="2,4,6",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_REORDER_PAGES,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_reorder",
            page_spec="3,1,2,2",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_WATERMARK_TEXT,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_watermark_text",
            watermark_text="CONFIDENTIAL",
            watermark_font_size=40,
            watermark_rotation=45,
            watermark_opacity=0.16,
            watermark_position="center",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_WATERMARK_IMAGE,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_watermark_image",
            watermark_image=sample["watermark"],
            watermark_opacity=0.28,
            watermark_position="bottom-right",
            watermark_image_scale_percent=28,
        ),
        PdfToolConfig(
            tool=PDF_TOOL_TEXT_OVERLAY,
            files=[sample["pdf"]],
            output_dir=Path("outputs") / "pdf_tool_text_overlay",
            page_spec="1",
            watermark_text="REVIEWED",
            watermark_font_size=24,
            watermark_rotation=0,
            watermark_opacity=0.92,
            watermark_position="top-right",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_IMAGE_OVERLAY,
            files=[sample["pdf"]],
            output_dir=Path("outputs") / "pdf_tool_image_overlay",
            page_spec="1",
            watermark_image=sample["watermark"],
            watermark_opacity=0.95,
            watermark_position="top-left",
            watermark_image_scale_percent=18,
        ),
        PdfToolConfig(
            tool=PDF_TOOL_REDACT_TEXT,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_redact_text",
            watermark_text="Page marker",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_REDACT_AREA,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_redact_area",
            page_spec="1",
            redact_rect="60,700,260,760",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_EDIT_TEXT,
            files=[sample["pdf_multi"]],
            output_dir=Path("outputs") / "pdf_tool_edit_text",
            page_spec="1",
            watermark_text="Data row 1A",
            replacement_text="Updated row 1A",
        ),
        PdfToolConfig(
            tool=PDF_TOOL_SIGN_VISIBLE,
            files=[sample["pdf"]],
            output_dir=Path("outputs") / "pdf_tool_sign_visible",
            watermark_text="Gokul Saraswat",
            watermark_image=sample["watermark"],
            watermark_opacity=0.98,
            watermark_position="bottom-right",
            watermark_image_scale_percent=16,
        ),
        PdfToolConfig(
            tool=PDF_TOOL_EDIT_METADATA,
            files=[sample["pdf"]],
            output_dir=Path("outputs") / "pdf_tool_edit_metadata",
            metadata_title="Patch 5 Sample Title",
            metadata_author="Gokul Omni Convert Lite",
            metadata_subject="Smoke test metadata",
            metadata_keywords="patch5,metadata,smoke",
            metadata_clear_existing=True,
        ),
    ]


def validate_advanced_outputs(tool_outputs: dict[str, list[Path]]) -> None:
    redacted_files = tool_outputs.get(PDF_TOOL_REDACT_TEXT, [])
    if not redacted_files:
        raise AssertionError("Redaction output was not created.")
    with fitz.open(str(redacted_files[0])) as document:
        text = "\n".join(page.get_text("text") for page in document)
    if "Page marker" in text:
        raise AssertionError("Redacted PDF still contains the redacted phrase.")

    area_redacted = tool_outputs.get(PDF_TOOL_REDACT_AREA, [])
    if not area_redacted:
        raise AssertionError("Area redaction output was not created.")
    with fitz.open(str(area_redacted[0])) as document:
        if document.page_count < 1:
            raise AssertionError("Area-redacted PDF is empty.")

    edit_text_files = tool_outputs.get(PDF_TOOL_EDIT_TEXT, [])
    if not edit_text_files:
        raise AssertionError("Best-effort text edit output was not created.")
    with fitz.open(str(edit_text_files[0])) as document:
        edited_text = "\n".join(page.get_text("text") for page in document)
    if "Updated row 1A" not in edited_text:
        raise AssertionError("Best-effort text edit did not insert the replacement text.")

    metadata_files = tool_outputs.get(PDF_TOOL_EDIT_METADATA, [])
    if not metadata_files:
        raise AssertionError("Metadata output was not created.")
    with fitz.open(str(metadata_files[0])) as document:
        metadata = document.metadata
    assert metadata.get("title") == "Patch 5 Sample Title"
    assert metadata.get("author") == "Gokul Omni Convert Lite"
    assert metadata.get("subject") == "Smoke test metadata"

    for required_tool in (PDF_TOOL_TEXT_OVERLAY, PDF_TOOL_IMAGE_OVERLAY, PDF_TOOL_SIGN_VISIBLE):
        if not tool_outputs.get(required_tool):
            raise AssertionError(f"Expected output for {required_tool} was not created.")


def validate_security_outputs(lock_output: Path, unlock_output: Path, compress_output: Path, original_pdf: Path) -> None:
    with fitz.open(str(lock_output)) as locked_doc:
        if not locked_doc.needs_pass:
            raise AssertionError("Locked PDF should require a password.")
        if locked_doc.authenticate("secret123") <= 0:
            raise AssertionError("Could not authenticate the locked PDF with the expected password.")
        if locked_doc.page_count < 1:
            raise AssertionError("Locked PDF did not expose pages after authentication.")

    with fitz.open(str(unlock_output)) as unlocked_doc:
        if unlocked_doc.needs_pass:
            raise AssertionError("Unlocked PDF should not remain encrypted.")
        if unlocked_doc.page_count < 1:
            raise AssertionError("Unlocked PDF should contain at least one page.")

    if not compress_output.exists() or compress_output.stat().st_size <= 0:
        raise AssertionError("Compressed PDF output was not created correctly.")

    with fitz.open(str(original_pdf)) as original_doc, fitz.open(str(compress_output)) as compressed_doc:
        if original_doc.page_count != compressed_doc.page_count:
            raise AssertionError("Compressed PDF page count changed unexpectedly.")


def validate_organizer_outputs(
    original_pdf: Path,
    organized_output: Path,
    extracted_output: Path,
    image_outputs: list[Path],
) -> None:
    summary = organizer_pdf_summary(original_pdf)
    if summary.page_count != 6:
        raise AssertionError("Organizer summary did not report the expected page count.")

    with fitz.open(str(organized_output)) as organized_doc:
        if organized_doc.page_count != 7:
            raise AssertionError("Organizer save output should contain seven pages after duplication.")
        first_page_text = organized_doc[0].get_text()
        second_page_text = organized_doc[1].get_text()
        third_page_text = organized_doc[2].get_text()
        if "Page marker: 2" not in first_page_text:
            raise AssertionError("Organizer ordering did not move the first page down as expected.")
        if "Page marker: 1" not in second_page_text or "Page marker: 1" not in third_page_text:
            raise AssertionError("Organizer duplication did not preserve the expected page text.")
        if organized_doc[2].rotation != 90:
            raise AssertionError("Organizer rotation did not persist in the saved PDF.")

    with fitz.open(str(extracted_output)) as extracted_doc:
        if extracted_doc.page_count != 3:
            raise AssertionError("Organizer extract output should contain the selected subset only.")
        extracted_first = extracted_doc[0].get_text()
        if "Page marker: 2" not in extracted_first:
            raise AssertionError("Organizer extract output did not follow the current organized order.")

    if len(image_outputs) != 2:
        raise AssertionError("Organizer image export should create exactly two images in the smoke test.")
    for image_path in image_outputs:
        if not image_path.exists() or image_path.stat().st_size <= 0:
            raise AssertionError(f"Organizer image export output is missing or empty: {image_path}")




def validate_patch19_organizer_layout_helpers() -> None:
    sequence = build_default_sequence(5)
    reordered, selected = move_positions_to_index(sequence, [0, 1], 4)
    expected = [2, 3, 0, 1, 4]
    if [page.source_index for page in reordered] != expected:
        raise AssertionError("Organizer drag-style reorder helper did not place the selected block at the expected slot.")
    if selected != [2, 3]:
        raise AssertionError("Organizer drag-style reorder helper did not return the updated selected positions.")

    rotated = rotate_positions(reordered, selected, 90)
    payload = sequence_to_payload(
        rotated,
        source_pdf="sample.pdf",
        page_count=5,
        selected_positions=selected,
    )
    restored_sequence, restored_selected = sequence_from_payload(payload, 5)
    if restored_sequence != rotated or restored_selected != selected:
        raise AssertionError("Organizer layout payload round-trip did not preserve the page sequence and selection.")


def validate_patch14_state_and_logs(outputs: Path) -> None:
    state_path = outputs / "state" / "app_state.json"
    store = AppStateStore(state_path)
    ensure_install_date(store.state)
    store.remember_outputs([outputs / "one.pdf", outputs / "two.pdf"])
    store.add_failed_job({"mode": "PDF Tool -> Compress", "job_type": "pdf_tool", "input_files": ["missing.pdf"]})
    store.set_session_snapshot({"mode": MODE_ANY_TO_PDF, "selected_files": ["a.pdf"]})
    store.update(
        auto_open_output_folder=True,
        restore_last_session=True,
        cleanup_temp_on_exit=True,
        update_checker_enabled=True,
        last_update_check="2026-04-15 10:00:00",
    )
    reloaded = AppStateStore(state_path)
    if len(reloaded.recent_outputs()) != 2:
        raise AssertionError("Recent outputs did not persist in AppStateStore.")
    if len(reloaded.failed_jobs()) != 1:
        raise AssertionError("Failed jobs did not persist in AppStateStore.")
    if reloaded.session_snapshot().get("mode") != MODE_ANY_TO_PDF:
        raise AssertionError("Session snapshot did not persist correctly.")
    if not reloaded.get("auto_open_output_folder") or not reloaded.get("restore_last_session") or not reloaded.get("cleanup_temp_on_exit"):
        raise AssertionError("Patch 14 settings flags did not persist correctly.")

    log_path = export_text_file(outputs / "logs" / "patch14_logs.txt", "Patch 14 log export smoke\n")
    if "Patch 14 log export smoke" not in log_path.read_text(encoding="utf-8"):
        raise AssertionError("Patch 14 text log export helper did not write the expected content.")


def run_link_download_smoke(inputs: Path, outputs: Path) -> tuple[list[Path], list[str]]:
    server = ThreadingHTTPServer(("127.0.0.1", 0), partial(SimpleHTTPRequestHandler, directory=str(inputs)))
    server_thread = threading.Thread(target=server.serve_forever, daemon=True)
    server_thread.start()
    cache_dir = cache_root_from_setting("", outputs)
    try:
        base_url = f"http://127.0.0.1:{server.server_port}"
        pasted = f"{base_url}/sample.pdf\n{base_url}/sample.md\n{base_url}/sample.html\n{base_url}/sample.pdf\n"
        extracted = extract_urls(pasted)
        if len(extracted) != 3:
            raise AssertionError("URL extraction did not deduplicate repeated links as expected.")
        results = download_many_urls(extracted, cache_dir, timeout=10)
        downloaded = [Path(item.file_path) for item in results if getattr(item, "status", "") == "downloaded" and getattr(item, "file_path", "")]
        if len(downloaded) != 3:
            raise AssertionError("Link downloader did not fetch the expected three local test assets.")
        converted_outputs: list[Path] = []
        converted_outputs.extend(
            process_batch(
                BatchConfig(mode=MODE_MD_TO_HTML, files=[downloaded[1]], output_dir=outputs / "links_md_to_html", engine_mode=ENGINE_PURE_PYTHON)
            )
        )
        converted_outputs.extend(
            process_batch(
                BatchConfig(mode=MODE_HTML_TO_PDF, files=[downloaded[2]], output_dir=outputs / "links_html_to_pdf", engine_mode=ENGINE_PURE_PYTHON)
            )
        )
        converted_outputs.extend(
            process_batch(
                BatchConfig(mode=MODE_PDF_TO_IMAGES, files=[downloaded[0]], output_dir=outputs / "links_pdf_to_images", image_format="png", image_scale=1.2)
            )
        )
        return downloaded + converted_outputs, []
    finally:
        server.shutdown()
        server.server_close()



def run_patch15_workflow_smoke(outputs: Path) -> list[Path]:
    cache_dir = outputs / "workflow_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)
    fresh = cache_dir / "fresh.bin"
    fresh.write_bytes(b"fresh-data")
    old = cache_dir / "old.bin"
    old.write_bytes(b"stale-data" * 200)
    old_time = (datetime.now() - timedelta(days=45)).timestamp()
    os.utime(old, (old_time, old_time))

    stats_before = directory_stats(cache_dir)
    if stats_before.file_count != 2:
        raise AssertionError("Patch 15 cache stats did not count the expected files.")

    prune_result = prune_directory(cache_dir, max_age_days=30, max_total_bytes=1024 * 1024)
    if int(prune_result.get("removed_count", 0)) < 1:
        raise AssertionError("Patch 15 cache prune did not remove the stale file.")
    if old.exists():
        raise AssertionError("Patch 15 cache prune left an expired file behind.")

    state_path = outputs / "state" / "patch15_app_state.json"
    store = AppStateStore(path=state_path)
    store.save_preset(
        {
            "name": "Patch15 Favorite",
            "mode": MODE_MD_TO_PDF,
            "output_dir": str(outputs / "favorite"),
            "engine_mode": ENGINE_PURE_PYTHON,
            "favorite": False,
        }
    )
    store.set_preset_favorite("Patch15 Favorite", True)
    favorites = store.favorite_presets()
    if not favorites or str(favorites[0].get("name", "")) != "Patch15 Favorite":
        raise AssertionError("Patch 15 preset favorite storage did not round-trip correctly.")

    report = outputs / "diagnostics" / "patch15_cache_summary.json"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(json.dumps({"before": stats_before.to_dict(), "after": prune_result}, indent=2), encoding="utf-8")
    return [fresh, report]





def run_patch16_release_smoke(outputs: Path) -> list[Path]:
    release_dir = outputs / "patch16_release"
    installer_dir = release_dir / "installer"
    installer_dir.mkdir(parents=True, exist_ok=True)

    state_path = release_dir / "app_state.json"
    state_path.write_text(json.dumps({"theme": "dark", "output_dir": str(outputs)}, indent=2), encoding="utf-8")
    notes_path = release_dir / "footer_notes.md"
    notes_path.write_text("# Patch 16 notes\n\nWorkspace bundle smoke test.\n", encoding="utf-8")
    about_path = release_dir / "about_profile.json"
    about_path.write_text(json.dumps({"name": "Gokul Saraswat", "image_path": "profile.png"}, indent=2), encoding="utf-8")
    image_path = release_dir / "profile.png"
    Image.new("RGB", (64, 64), color=(32, 80, 140)).save(image_path)
    static_about_path = installer_dir / "about_static.json"
    static_about_path.write_text(json.dumps({"name": "Static About"}, indent=2), encoding="utf-8")
    asset_config_path = release_dir / "remote_assets.json"
    save_asset_config({"remote_enabled": False}, asset_config_path)
    build_notes = installer_dir / "BUILDING.md"
    build_notes.write_text("# Build notes\n\nPatch 16 release workspace smoke.\n", encoding="utf-8")

    manifest_path = build_example_update_manifest(installer_dir / "update_manifest.example.json", "1.8.0 Patch 18")
    manifest_payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest_payload["version"] = "1.8.1 Patch 18"
    manifest_payload["notes"] = "Patch 16 smoke manifest."
    manifest_path.write_text(json.dumps(manifest_payload, indent=2), encoding="utf-8")

    update_result = check_for_updates("1.8.0 Patch 18", str(manifest_path))
    if not update_result.get("has_update"):
        raise AssertionError(f"Patch 16 update checker did not detect the newer manifest: {update_result}")

    bundle_path = export_workspace_bundle(
        release_dir / "workspace_bundle.zip",
        state_path=state_path,
        notes_path=notes_path,
        about_profile_path=about_path,
        static_about_profile_path=static_about_path,
        installer_dir=installer_dir,
        asset_config_path=asset_config_path,
        extra_files=[image_path, manifest_path],
    )
    if not bundle_path.exists():
        raise AssertionError("Patch 16 workspace bundle export did not create the ZIP file.")

    import_root = release_dir / "imported_bundle"
    summary = import_workspace_bundle(bundle_path, import_root)
    extracted = summary.get("extracted", [])
    if int(summary.get("extracted_count", 0)) < 5 or not extracted:
        raise AssertionError(f"Patch 16 workspace bundle import extracted too few files: {summary}")
    if not (import_root / "workspace_manifest.json").exists():
        raise AssertionError("Patch 16 workspace import did not include the bundle manifest file.")

    return [manifest_path, bundle_path, import_root / "workspace_manifest.json"]



def run_patch18_accessibility_smoke(outputs: Path) -> list[Path]:
    patch_dir = outputs / "patch18_accessibility"
    patch_dir.mkdir(parents=True, exist_ok=True)
    state_path = patch_dir / "app_state.json"
    store = AppStateStore(path=state_path)
    store.update(
        theme="dark",
        compact_ui=True,
        ui_scale="125%",
        high_contrast=True,
        reduced_motion=True,
        state_backup_enabled=True,
        state_backup_keep=6,
        output_dir=str(outputs / "patch18_scaled"),
    )
    store.update(start_page="settings", last_page="settings", support_bundle_dir=str(outputs / "support"))
    store.update(last_update_check="2026-04-15 11:00:00")
    backup_path = latest_state_backup(state_path)
    if backup_path is None or not backup_path.exists():
        raise AssertionError("Patch 18 automatic state backup did not create a backup file.")
    if store.get("ui_scale") != "125%" or not store.get("high_contrast") or not store.get("reduced_motion"):
        raise AssertionError("Patch 18 accessibility flags did not persist in AppStateStore.")
    backup_root = state_backup_dir(state_path)
    if not backup_root.exists():
        raise AssertionError("Patch 18 backup directory helper did not resolve an existing folder.")

    state_path.write_text("{broken json", encoding="utf-8")
    recovered = AppStateStore(path=state_path)
    if recovered.get("ui_scale") != "125%" or not recovered.get("high_contrast"):
        raise AssertionError("Patch 18 backup recovery did not restore the latest valid state.")
    if str(recovered.get("start_page")) != "settings":
        raise AssertionError("Patch 18 backup recovery lost the persisted start page.")

    shortcuts_path = Path(__file__).with_name("keyboard_shortcuts.md")
    if not shortcuts_path.exists():
        raise AssertionError("Patch 18 shortcut guide file is missing.")
    shortcuts_text = shortcuts_path.read_text(encoding="utf-8")
    if "Ctrl+Enter" not in shortcuts_text or "F1" not in shortcuts_text:
        raise AssertionError("Patch 18 shortcut guide did not contain the expected default shortcuts.")

    backup_report = patch_dir / "backup_report.json"
    backup_report.write_text(
        json.dumps(
            {
                "backup_dir": str(backup_root),
                "latest_backup": str(backup_path),
                "restored_theme": recovered.get("theme"),
                "ui_scale": recovered.get("ui_scale"),
                "high_contrast": recovered.get("high_contrast"),
                "reduced_motion": recovered.get("reduced_motion"),
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    return [backup_path, backup_report, shortcuts_path]





def run_patch23_asset_smoke(inputs: Path, outputs: Path, sample: dict[str, Path]) -> None:
    remote_header = inputs / "remote_header.gif"
    Image.open(sample["img1"]).save(remote_header, format="GIF")

    remote_profile = inputs / "remote_about.json"
    remote_profile.write_text(
        json.dumps(
            {
                "name": "Remote Gokul",
                "title": "Remote Profile",
                "subtitle": "Loaded from a test URL",
                "company": "Oracle Corporation",
                "project": APP_NAME,
                "email": "gokul.saraswat@oracle.com",
                "handle": "@gokul.saraswat",
                "bio": "Remote profile smoke test.",
                "image_path": "assets/gokul_profile_placeholder.png",
                "image_url": "",
                "feedback_url": "mailto:gokul.saraswat@oracle.com",
                "contribute_url": "https://github.com/gokul-saraswat/gokul-omni-convert-lite/issues",
                "links": [],
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    config_path = outputs / "remote_assets.json"
    config = load_asset_config(config_path)
    config.update(
        {
            "remote_enabled": True,
            "cache_dir": str(outputs / "asset_cache"),
            "timeout": 10,
            "refresh_hours": 24,
        }
    )

    server = ThreadingHTTPServer(("127.0.0.1", 0), partial(SimpleHTTPRequestHandler, directory=str(inputs)))
    server_thread = threading.Thread(target=server.serve_forever, daemon=True)
    server_thread.start()
    try:
        base_url = f"http://127.0.0.1:{server.server_port}"
        config["header_gif_url"] = f"{base_url}/remote_header.gif"
        config["splash_gif_url"] = f"{base_url}/remote_header.gif"
        config["profile_json_url"] = f"{base_url}/remote_about.json"
        save_asset_config(config, config_path)
        loaded = load_asset_config(config_path)

        header_info = resolve_local_or_remote_asset(
            "assets/gokul_header.gif",
            str(loaded.get("header_gif_url", "")),
            base_dir=Path(__file__).resolve().parent,
            fallback_value=Path("assets") / "gokul_header.gif",
            config=loaded,
        )
        header_path = Path(str(header_info.get("path", "")))
        if not header_path.exists():
            raise AssertionError("Patch 23 remote header asset did not resolve to an existing cached file.")

        pulled_profile = outputs / "pulled_about.json"
        download_text_file(str(loaded.get("profile_json_url", "")), pulled_profile, timeout=10)
        pulled_data = json.loads(pulled_profile.read_text(encoding="utf-8"))
        if pulled_data.get("name") != "Remote Gokul":
            raise AssertionError("Patch 23 remote profile pull did not download the expected JSON.")

        summary = asset_cache_summary(loaded)
        if int(summary.get("count", 0)) < 1:
            raise AssertionError("Patch 23 asset cache summary did not record downloaded remote assets.")

        clear_asset_cache(loaded)
        cleared = asset_cache_summary(loaded)
        if int(cleared.get("count", 0)) != 0:
            raise AssertionError("Patch 23 asset cache clear did not remove cached assets.")
    finally:
        server.shutdown()
        server.server_close()


def run() -> None:
    with tempfile.TemporaryDirectory() as temp_dir:
        root = Path(temp_dir)
        inputs = root / "inputs"
        outputs = root / "outputs"
        inputs.mkdir()
        outputs.mkdir()
        sample = create_sample_files(inputs)

        profile = load_about_profile(Path(__file__).with_name("about_profile.json"))
        validate_ui_text_helpers()
        validate_flow_wrap_helper()
        validate_preview_support(sample)
        if not profile.get("feedback_url"):
            raise AssertionError("Patch 12 About profile did not expose a feedback URL.")
        if not profile.get("company"):
            raise AssertionError("Patch 12 About profile did not expose a company field.")

        run_patch23_asset_smoke(inputs, outputs, sample)

        splash_state = {
            "splash_enabled": True,
            "splash_seen": False,
        }
        if not should_show_first_launch_splash(splash_state):
            raise AssertionError("Patch 12 splash logic did not detect the first launch.")
        splash_state["splash_seen"] = True
        if should_show_first_launch_splash(splash_state):
            raise AssertionError("Patch 12 splash logic did not stop after the first launch.")

        login_state = {
            "install_date": (datetime.now() - timedelta(days=4)).isoformat(timespec="seconds"),
            "login_popup_enabled": True,
            "login_popup_dismissed": False,
            "login_popup_completed": False,
            "login_popup_last_shown": "",
        }
        if not should_show_login_popup(dict(login_state)):
            raise AssertionError("Patch 12 login reminder did not become eligible after 3+ days.")
        login_state["login_popup_dismissed"] = True
        if should_show_login_popup(dict(login_state)):
            raise AssertionError("Patch 12 login reminder ignored the dismissed state.")
        login_state["login_popup_dismissed"] = False
        login_state["login_popup_completed"] = True
        if should_show_login_popup(dict(login_state)):
            raise AssertionError("Patch 12 login reminder ignored the completed state.")
        login_state["login_popup_completed"] = False
        login_state["login_popup_last_shown"] = datetime.now().isoformat(timespec="seconds")
        if should_show_login_popup(dict(login_state)):
            raise AssertionError("Patch 12 login reminder ignored the cooldown window.")

        conversion_jobs, skipped = build_conversion_jobs(sample)
        pdf_tool_jobs = build_pdf_tool_jobs(sample)

        all_outputs: list[Path] = []
        tool_outputs: dict[str, list[Path]] = {}

        for job in conversion_jobs:
            job.output_dir = outputs / job.output_dir.name
            all_outputs.extend(process_batch(job))

        for job in pdf_tool_jobs:
            job.output_dir = outputs / job.output_dir.name
            created = process_pdf_tool(job)
            tool_outputs[job.tool] = created
            all_outputs.extend(created)

        lock_job = PdfToolConfig(
            tool=PDF_TOOL_LOCK,
            files=[sample["pdf"]],
            output_dir=outputs / "pdf_tool_lock",
            pdf_password="secret123",
            pdf_owner_password="owner123",
        )
        locked_outputs = process_pdf_tool(lock_job)
        tool_outputs[lock_job.tool] = locked_outputs
        all_outputs.extend(locked_outputs)

        unlock_job = PdfToolConfig(
            tool=PDF_TOOL_UNLOCK,
            files=locked_outputs,
            output_dir=outputs / "pdf_tool_unlock",
            pdf_password="secret123",
        )
        unlocked_outputs = process_pdf_tool(unlock_job)
        tool_outputs[unlock_job.tool] = unlocked_outputs
        all_outputs.extend(unlocked_outputs)

        compress_job = PdfToolConfig(
            tool=PDF_TOOL_COMPRESS,
            files=[sample["pdf_multi"]],
            output_dir=outputs / "pdf_tool_compress",
            compression_profile="balanced",
        )
        compressed_outputs = process_pdf_tool(compress_job)
        tool_outputs[compress_job.tool] = compressed_outputs
        all_outputs.extend(compressed_outputs)

        organizer_sequence = build_default_sequence(6)
        organizer_sequence, moved_positions = move_positions_down(organizer_sequence, [0])
        organizer_sequence, duplicated_positions = duplicate_positions(organizer_sequence, moved_positions)
        organizer_sequence = rotate_positions(organizer_sequence, duplicated_positions, 90)

        organizer_save_output = outputs / "organizer_save" / "sample_multi_organized.pdf"
        organizer_save_output.parent.mkdir(parents=True, exist_ok=True)
        organizer_save_sequence_as_pdf(sample["pdf_multi"], organizer_sequence, organizer_save_output)

        organizer_extract_output = outputs / "organizer_extract" / "sample_multi_selected.pdf"
        organizer_extract_output.parent.mkdir(parents=True, exist_ok=True)
        organizer_extract_selected_pdf(sample["pdf_multi"], organizer_sequence, [0, 2, 6], organizer_extract_output)

        organizer_image_output_dir = outputs / "organizer_export_images"
        organizer_image_outputs = organizer_export_pages_as_images(
            sample["pdf_multi"],
            organizer_sequence,
            [0, 2],
            organizer_image_output_dir,
        )
        all_outputs.extend([organizer_save_output, organizer_extract_output, *organizer_image_outputs])

        validate_advanced_outputs(tool_outputs)
        validate_security_outputs(locked_outputs[0], unlocked_outputs[0], compressed_outputs[0], sample["pdf_multi"])
        validate_organizer_outputs(sample["pdf_multi"], organizer_save_output, organizer_extract_output, organizer_image_outputs)
        validate_patch19_organizer_layout_helpers()

        html_output = outputs / "pdf_to_html" / "sample_multi.html"
        if not html_output.exists():
            raise AssertionError("PDF -> HTML output was not created.")
        html_text = html_output.read_text(encoding="utf-8")
        if "Page 1" not in html_text or "Sample multi-page PDF page 1" not in html_text:
            raise AssertionError("PDF -> HTML output did not contain expected page text.")

        md_html_output = outputs / "md_to_html" / "sample.html"
        if not md_html_output.exists():
            raise AssertionError("Markdown -> HTML output was not created.")

        pptx_output = outputs / "pdf_to_pptx" / "sample.pptx"
        if not pptx_output.exists():
            raise AssertionError("PDF -> PPTX output was not created.")
        try:
            from pptx import Presentation

            presentation = Presentation(str(pptx_output))
            if len(presentation.slides) != 1:
                raise AssertionError("PDF -> PPTX output should contain one slide for the one-page sample PDF.")
        except Exception as exc:
            raise AssertionError(f"Could not validate PDF -> PPTX output: {exc}") from exc

        if sample.get("pptx"):
            presentation_image_dir = outputs / "presentations_to_images" / "sample_presentation"
            image_candidates = sorted(presentation_image_dir.glob("*.png"))
            if not image_candidates:
                raise AssertionError("Presentation -> Images output was not created.")

        html_docx_output = outputs / "html_to_docx" / "sample.docx"
        if not html_docx_output.exists():
            raise AssertionError("HTML -> DOCX output was not created.")

        html_md_output = outputs / "html_to_md" / "sample.md"
        if not html_md_output.exists() or "Sample HTML" not in html_md_output.read_text(encoding="utf-8"):
            raise AssertionError("HTML -> Markdown output was not created correctly.")

        explicit_md_pdf = outputs / "md_to_pdf_explicit" / "sample.pdf"
        explicit_html_pdf = outputs / "html_to_pdf_explicit" / "sample.pdf"
        if not explicit_md_pdf.exists() or not explicit_html_pdf.exists():
            raise AssertionError("Explicit Markdown/HTML PDF outputs were not created.")

        pure_docx_pdf = outputs / "docx_to_pdf_pure_python" / "sample.pdf"
        pure_sheet_pdf = outputs / "sheets_to_pdf_pure_python" / "sample.pdf"
        pure_md_pdf = outputs / "markdown_to_pdf_pure_python" / "sample.pdf"
        pure_html_pdf = outputs / "html_to_pdf_pure_python" / "sample.pdf"
        if not pure_docx_pdf.exists() or not pure_sheet_pdf.exists() or not pure_md_pdf.exists() or not pure_html_pdf.exists():
            raise AssertionError("Pure Python Patch 6 outputs were not created.")

        with fitz.open(str(pure_docx_pdf)) as docx_pdf_doc:
            extracted = "\n".join(page.get_text() for page in docx_pdf_doc)
            if "Header A" not in extracted or "Second bullet item" not in extracted:
                raise AssertionError("Pure Python DOCX PDF did not include the expected rich content.")

        with fitz.open(str(pure_sheet_pdf)) as sheet_pdf_doc:
            extracted = "\n".join(page.get_text() for page in sheet_pdf_doc)
            if "Summary" not in extracted or "Metric" not in extracted:
                raise AssertionError("Pure Python spreadsheet PDF did not include the expected sheet structure.")

        with fitz.open(str(pure_md_pdf)) as md_pdf_doc:
            extracted = "\n".join(page.get_text() for page in md_pdf_doc)
            if "Sample Markdown" not in extracted or "quoted line" not in extracted:
                raise AssertionError("Pure Python Markdown PDF did not include expected structured content.")

        with fitz.open(str(pure_html_pdf)) as html_pdf_doc:
            extracted = "\n".join(page.get_text() for page in html_pdf_doc)
            if "Sample HTML" not in extracted or "Row A" not in extracted:
                raise AssertionError("Pure Python HTML PDF did not include expected structured content.")

        route_preview = build_conversion_route_preview(MODE_TEXT_TO_PDF, [sample["md"], sample["html"]], engine_mode=ENGINE_PURE_PYTHON)
        if "structured markdown" not in route_preview.lower() or "structure-aware" not in route_preview.lower():
            raise AssertionError("Route preview did not expose the expected Patch 6 fidelity labels.")
        route_preview_html = build_conversion_route_preview(MODE_HTML_TO_PDF, [sample["html"]], engine_mode=ENGINE_PURE_PYTHON)
        if "html" not in route_preview_html.lower():
            raise AssertionError("Patch 13 route preview did not expose the HTML pipeline.")

        watch_root = root / "watch_inputs"
        watch_root.mkdir(parents=True, exist_ok=True)
        watch_txt = watch_root / "watch_sample.txt"
        watch_md = watch_root / "watch_sample.md"
        watch_txt.write_text("watch text", encoding="utf-8")
        watch_md.write_text("# watch markdown", encoding="utf-8")
        watch_files, watch_fingerprints = discover_watch_candidates(watch_root, MODE_TEXT_TO_PDF, True, [])
        if len(watch_files) != 2 or len(watch_fingerprints) != 2:
            raise AssertionError("Automation watch discovery did not find the expected files.")
        seen = add_fingerprints([], watch_fingerprints)
        repeat_files, _ = discover_watch_candidates(watch_root, MODE_TEXT_TO_PDF, True, seen)
        if repeat_files:
            raise AssertionError("Automation watch discovery should skip already-seen files.")

        report_path = write_run_report(
            {
                "timestamp": "2026-04-14 12:00:00",
                "status": "Completed",
                "mode": MODE_TEXT_TO_PDF,
                "job_type": "convert",
                "file_count": 2,
                "output_count": 1,
                "output_dir": str(outputs / "automation_report"),
                "inputs_preview": [str(watch_txt), str(watch_md)],
                "outputs_preview": [str(pure_md_pdf)],
            },
            outputs / "automation_report" / "last_run_report.txt",
        )
        if not report_path.exists() or "Gokul Omni Convert Lite Run Report" not in report_path.read_text(encoding="utf-8"):
            raise AssertionError("Automation report helper did not create the expected report file.")

        bundle_path = bundle_paths_as_zip([pure_md_pdf, pure_html_pdf, report_path], outputs / "automation_bundle" / "outputs_bundle.zip")
        if not bundle_path.exists():
            raise AssertionError("Automation ZIP bundle helper did not create the bundle.")

        preset_export = outputs / "automation_presets" / "presets.json"
        export_presets_to_json(
            [
                {
                    "name": "Docs to PDF",
                    "mode": MODE_DOCS_TO_PDF,
                    "engine_mode": ENGINE_PURE_PYTHON,
                    "merge_to_one_pdf": False,
                },
                {
                    "name": "Markdown bundle",
                    "mode": MODE_TEXT_TO_PDF,
                    "engine_mode": ENGINE_PURE_PYTHON,
                    "merge_to_one_pdf": True,
                },
            ],
            preset_export,
        )
        imported_presets = import_presets_from_json(preset_export)
        if len(imported_presets) != 2 or {item["name"] for item in imported_presets} != {"Docs to PDF", "Markdown bundle"}:
            raise AssertionError("Preset export/import helpers did not round-trip correctly.")

        state_snapshot = export_state_snapshot(
            outputs / "state_snapshot" / "settings_snapshot.json",
            {
                "theme": "dark",
                "output_dir": str(outputs / "state_snapshot_output"),
                "smtp_settings": SMTPSettings(host="smtp.example.com", port=587, sender="sender@example.com").to_state_dict(),
            },
        )
        imported_state = import_state_snapshot(state_snapshot)
        if imported_state.get("theme") != "dark" or "smtp_settings" not in imported_state:
            raise AssertionError("State snapshot export/import helpers did not round-trip correctly.")

        validate_patch14_state_and_logs(outputs)

        link_outputs, link_skips = run_link_download_smoke(inputs, outputs)
        all_outputs.extend(link_outputs)
        skipped.extend(link_skips)

        patch15_outputs = run_patch15_workflow_smoke(outputs)
        all_outputs.extend(patch15_outputs)

        patch16_outputs = run_patch16_release_smoke(outputs)
        all_outputs.extend(patch16_outputs)

        patch18_outputs = run_patch18_accessibility_smoke(outputs)
        all_outputs.extend(patch18_outputs)

        diagnostics_path = export_diagnostics_report(
            outputs / "diagnostics" / "diagnostics.json",
            app_name=APP_NAME,
            app_version="1.8.0 Patch 18",
            state_path=APP_STATE_PATH,
            about_profile_path=Path("about_profile.json"),
            notes_path=Path("footer_notes.md"),
            installer_dir=Path("installer"),
            asset_config_path=Path("remote_assets.json"),
            output_dir=outputs,
            selected_files=[str(sample["docx"])],
            last_outputs=[str(pure_md_pdf), str(pure_html_pdf)],
            dependency_status=dependency_status(),
            smtp_summary=SMTPSettings(host="smtp.example.com", port=587, sender="sender@example.com").sanitized_dict(),
            extra={"smoke": True},
        )
        diagnostics_text = diagnostics_path.read_text(encoding="utf-8")
        if "installer_assets" not in diagnostics_text or "smtp" not in diagnostics_text:
            raise AssertionError("Diagnostics export did not include the expected Patch 9 sections.")

        activity_report = render_activity_report_html(
            outputs / "reports" / "activity_report.html",
            app_name=APP_NAME,
            app_version="1.8.0 Patch 18",
            recent_jobs=[
                {
                    "timestamp": "2026-04-15 11:30:00",
                    "status": "Completed",
                    "mode": MODE_TEXT_TO_PDF,
                    "file_count": 2,
                    "output_count": 2,
                    "output_dir": str(outputs),
                }
            ],
            recent_outputs=[str(pure_md_pdf), str(pure_html_pdf)],
            failed_jobs=[],
            dependency_status=dependency_status(),
            notes="Patch 17 smoke report.",
        )
        if not activity_report.exists() or "Activity Report" not in activity_report.read_text(encoding="utf-8"):
            raise AssertionError("Patch 17 activity report export did not create the expected HTML output.")

        support_bundle = export_support_bundle(
            outputs / "reports" / "support_bundle.zip",
            diagnostics_report=diagnostics_path,
            state_snapshot=state_snapshot,
            activity_report=activity_report,
            logs_path=report_path,
            notes_path=Path(__file__).with_name("footer_notes.md"),
            about_profile_path=Path(__file__).with_name("about_profile.json"),
            installer_dir=Path(__file__).with_name("installer"),
            asset_config_path=Path(__file__).with_name("remote_assets.json"),
            extra_files=[Path(__file__).with_name("README.md")],
        )
        if not support_bundle.exists():
            raise AssertionError("Patch 17 support bundle export did not create the ZIP file.")
        with zipfile.ZipFile(support_bundle) as archive:
            names = set(archive.namelist())
            required = {"manifest.json", "reports/activity_report.html", "reports/diagnostics.json", "reports/state_snapshot.json"}
            if not required.issubset(names):
                raise AssertionError(f"Patch 17 support bundle is missing expected files: {required - names}")

        message = build_email_message(
            sender="sender@example.com",
            recipients=["receiver@example.com"],
            subject="Patch 15 smoke",
            body="Testing attachment generation.",
            attachments=[pure_md_pdf],
        )
        if "receiver@example.com" not in str(message["To"]) or not message.is_multipart():
            raise AssertionError("Patch 11 email message builder did not create the expected MIME message.")

        eml_draft = build_eml_draft(
            outputs / "mail" / "latest_outputs.eml",
            sender="sender@example.com",
            recipients=["receiver@example.com"],
            subject="Patch 15 EML",
            body="Testing EML draft generation.",
            attachments=[pure_md_pdf, pure_html_pdf],
        )
        eml_text = eml_draft.read_text(encoding="utf-8", errors="replace")
        if "latest_outputs.eml" == eml_draft.name and "sample.pdf" not in eml_text:
            raise AssertionError("Patch 11 EML draft did not include the expected attachment references.")

        mailto_url = create_mailto_url(["receiver@example.com"], subject="Patch15", body="Hello from Patch 15")
        if not mailto_url.startswith("mailto:"):
            raise AssertionError("Patch 11 mailto helper did not build a mailto URL.")

        tesseract = detect_tesseract_status()
        if tesseract.get("available"):
            ocr_cfg = OcrConfig(language="eng", dpi=220, psm=6)
            ocr_output_dir = outputs / "ocr"
            try:
                image_searchable = image_to_searchable_pdf(sample["ocr_image"], ocr_output_dir / "ocr_image_searchable.pdf", config=ocr_cfg)
                pdf_searchable = pdf_to_searchable_pdf(sample["ocr_pdf"], ocr_output_dir / "ocr_pdf_searchable.pdf", config=ocr_cfg)
                ocr_text_path = extract_text_with_ocr(sample["ocr_image"], ocr_output_dir / "ocr_image.txt", config=ocr_cfg)
            except OcrError as exc:
                skipped.append(f"OCR searchable-PDF smoke tests were skipped: {exc}")
            else:
                for path in (image_searchable, pdf_searchable, ocr_text_path):
                    if not path.exists():
                        raise AssertionError(f"OCR output was not created: {path}")
                ocr_text = ocr_text_path.read_text(encoding="utf-8", errors="replace").upper()
                if "INVOICE" not in ocr_text:
                    raise AssertionError(f"OCR text extraction did not recover the expected text: {ocr_text}")
                with fitz.open(str(image_searchable)) as searchable_doc:
                    searchable_text = "\n".join(page.get_text("text") for page in searchable_doc).upper()
                if "INVOICE" not in searchable_text:
                    raise AssertionError("Image -> searchable PDF did not contain searchable OCR text.")
                all_outputs.extend([image_searchable, pdf_searchable, ocr_text_path])
        else:
            skipped.append("Tesseract not found. OCR smoke tests were skipped.")

        header_gif = Path(__file__).with_name("assets") / "gokul_header.gif"
        if not header_gif.exists():
            raise AssertionError("Patch 22 header GIF placeholder is missing.")
        print("Patch 32 smoke test completed successfully.")
        for note in skipped:
            print(f"SKIP: {note}")
        for path in all_outputs:
            print(path)

        mirror = Path.cwd() / "smoke_test_output_example"
        if mirror.exists():
            shutil.rmtree(mirror)
        shutil.copytree(outputs, mirror)
        print(f"Example outputs copied to: {mirror}")


if __name__ == "__main__":
    run()
