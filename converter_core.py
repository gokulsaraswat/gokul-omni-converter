from __future__ import annotations

import csv
import importlib.util
import os
import re
import shutil
import subprocess
import tempfile
import textwrap
import zipfile
from datetime import datetime
from html import escape
from html.parser import HTMLParser
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Callable, Iterable, Sequence
from xml.etree import ElementTree as ET

import fitz  # PyMuPDF
import pdfplumber
from PIL import Image
from docx import Document
from docx.shared import Pt
from openpyxl import Workbook, load_workbook
try:
    from pypdf import PdfReader, PdfWriter
except Exception:  # pragma: no cover - compatibility fallback for environments with PyPDF2 only
    from PyPDF2 import PdfReader, PdfWriter  # type: ignore
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from pure_python_renderers import (
    render_docx_to_pdf,
    render_html_to_docx_from_string,
    render_html_to_pdf_from_string,
    render_presentation_to_pdf,
    render_spreadsheet_to_pdf,
)


class ConversionError(RuntimeError):
    """Raised when a conversion cannot be completed."""


LogFn = Callable[[str], None]
ProgressFn = Callable[[int, int], None]

IMAGE_EXTS = {
    ".jpg",
    ".jpeg",
    ".png",
    ".bmp",
    ".tif",
    ".tiff",
    ".webp",
    ".gif",
}
DOC_EXTS = {".doc", ".docx", ".odt", ".rtf", ".fodt"}
SHEET_EXTS = {".xls", ".xlsx", ".ods", ".csv", ".tsv"}
TEXT_EXTS = {
    ".txt",
    ".text",
    ".log",
    ".ini",
    ".cfg",
    ".conf",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".toml",
}
MARKDOWN_EXTS = {".md", ".markdown", ".mdown"}
WEB_EXTS = {".html", ".htm"}
PRESENTATION_EXTS = {".ppt", ".pptx", ".odp"}
PDF_EXTS = {".pdf"}

AUTO_TO_PDF_EXTS = (
    IMAGE_EXTS
    | DOC_EXTS
    | SHEET_EXTS
    | TEXT_EXTS
    | MARKDOWN_EXTS
    | WEB_EXTS
    | PRESENTATION_EXTS
    | PDF_EXTS
)
TEXTUAL_TO_PDF_EXTS = TEXT_EXTS | MARKDOWN_EXTS | WEB_EXTS

MODE_ANY_TO_PDF = "Any Supported -> PDF"
MODE_IMAGES_TO_PDF = "Images -> PDF"
MODE_PDF_TO_IMAGES = "PDF -> Images"
MODE_DOCS_TO_PDF = "DOC / DOCX / ODT / RTF -> PDF"
MODE_PDF_TO_DOCX = "PDF -> DOCX"
MODE_SHEETS_TO_PDF = "XLS / XLSX / ODS / CSV -> PDF"
MODE_PDF_TO_XLSX = "PDF -> XLSX (tables/text)"
MODE_TEXT_TO_PDF = "Text / Markdown / HTML -> PDF"
MODE_MD_TO_PDF = "Markdown -> PDF"
MODE_MD_TO_DOCX = "Markdown -> DOCX"
MODE_MD_TO_HTML = "Markdown -> HTML"
MODE_HTML_TO_PDF = "HTML -> PDF"
MODE_HTML_TO_DOCX = "HTML -> DOCX"
MODE_HTML_TO_MD = "HTML -> Markdown"
MODE_MERGE_PDFS = "Merge PDFs"
MODE_PRESENTATIONS_TO_PDF = "PPT / PPTX / ODP -> PDF"
MODE_PRESENTATIONS_TO_IMAGES = "PPT / PPTX / ODP -> Images"
MODE_PDF_TO_PPTX = "PDF -> PPTX"
MODE_PDF_TO_HTML = "PDF -> HTML"

ENGINE_AUTO = "auto"
ENGINE_PURE_PYTHON = "pure_python"
ENGINE_LIBREOFFICE = "libreoffice"
ENGINE_ORDER = [ENGINE_AUTO, ENGINE_PURE_PYTHON, ENGINE_LIBREOFFICE]
ENGINE_HELP: dict[str, str] = {
    ENGINE_AUTO: (
        "Auto uses the richer built-in pure Python route first for DOCX, XLSX, PPTX, HTML, Markdown, TXT, CSV, and similar files, then falls back to LibreOffice only when needed."
    ),
    ENGINE_PURE_PYTHON: (
        "Pure Python keeps LibreOffice optional. It now uses structure-aware renderers for DOCX, XLSX, PPTX, HTML, and Markdown, with text-first fallback for legacy or partially supported formats."
    ),
    ENGINE_LIBREOFFICE: (
        "LibreOffice uses soffice for higher-fidelity Office-style rendering. Configure the soffice path in Settings or make it available in PATH when you want the external fallback route."
    ),
}

PDF_OUTPUT_MODES = {
    MODE_ANY_TO_PDF,
    MODE_IMAGES_TO_PDF,
    MODE_DOCS_TO_PDF,
    MODE_SHEETS_TO_PDF,
    MODE_TEXT_TO_PDF,
    MODE_MD_TO_PDF,
    MODE_HTML_TO_PDF,
    MODE_PRESENTATIONS_TO_PDF,
}

MODE_TO_EXTENSIONS: dict[str, set[str]] = {
    MODE_ANY_TO_PDF: AUTO_TO_PDF_EXTS,
    MODE_IMAGES_TO_PDF: IMAGE_EXTS,
    MODE_PDF_TO_IMAGES: PDF_EXTS,
    MODE_DOCS_TO_PDF: DOC_EXTS,
    MODE_PDF_TO_DOCX: PDF_EXTS,
    MODE_SHEETS_TO_PDF: SHEET_EXTS,
    MODE_PDF_TO_XLSX: PDF_EXTS,
    MODE_TEXT_TO_PDF: TEXTUAL_TO_PDF_EXTS,
    MODE_MD_TO_PDF: MARKDOWN_EXTS,
    MODE_MD_TO_DOCX: MARKDOWN_EXTS,
    MODE_MD_TO_HTML: MARKDOWN_EXTS,
    MODE_HTML_TO_PDF: WEB_EXTS,
    MODE_HTML_TO_DOCX: WEB_EXTS,
    MODE_HTML_TO_MD: WEB_EXTS,
    MODE_MERGE_PDFS: PDF_EXTS,
    MODE_PRESENTATIONS_TO_PDF: PRESENTATION_EXTS,
    MODE_PRESENTATIONS_TO_IMAGES: PRESENTATION_EXTS,
    MODE_PDF_TO_PPTX: PDF_EXTS,
    MODE_PDF_TO_HTML: PDF_EXTS,
}

MODE_HELP: dict[str, str] = {
    MODE_ANY_TO_PDF: (
        "Convert mixed supported files to PDF. You can batch convert or merge many sources "
        "into one final PDF. Supports images, Office files, spreadsheets, text files, "
        "HTML, Markdown, presentations, and existing PDFs."
    ),
    MODE_IMAGES_TO_PDF: (
        "Convert one or many images into PDFs. Enable the merge option to combine many images "
        "into a single PDF."
    ),
    MODE_PDF_TO_IMAGES: (
        "Export every page of each PDF as PNG or JPG images. Each input PDF gets its own output folder."
    ),
    MODE_DOCS_TO_PDF: (
        "Convert DOC, DOCX, ODT, or RTF files to PDF. You can convert each file separately or merge "
        "all results into one PDF."
    ),
    MODE_PDF_TO_DOCX: (
        "Best-effort PDF to DOCX extraction. Works best for text-heavy PDFs. Complex layouts and images "
        "may not round-trip perfectly."
    ),
    MODE_SHEETS_TO_PDF: (
        "Convert XLS, XLSX, ODS, CSV, or TSV files to PDF. You can convert each file separately or merge "
        "all resulting PDFs into one."
    ),
    MODE_PDF_TO_XLSX: (
        "Best-effort PDF table and text extraction into XLSX. Table-like PDFs work best. If tables are not "
        "detected, the app falls back to page text extraction."
    ),
    MODE_TEXT_TO_PDF: (
        "Convert TXT, Markdown, HTML, and similar text-like files to PDF. Patch 13 also adds dedicated Markdown and HTML conversion modes for faster discovery in the UI."
    ),
    MODE_MD_TO_PDF: (
        "Convert Markdown files directly to PDF using the built-in Markdown-to-HTML renderer and PDF export pipeline."
    ),
    MODE_MD_TO_DOCX: (
        "Convert Markdown files to DOCX. Pandoc is used when available, with a built-in fallback renderer for basic Markdown."
    ),
    MODE_MD_TO_HTML: (
        "Convert Markdown files to HTML. Pandoc is used when available, with a built-in fallback renderer for basic Markdown."
    ),
    MODE_HTML_TO_PDF: (
        "Convert HTML files directly to PDF. Pure Python HTML rendering is used first, with optional LibreOffice fallback only when your engine settings allow it."
    ),
    MODE_HTML_TO_DOCX: (
        "Convert HTML files to DOCX. Pandoc is used when available, with a built-in structured HTML fallback and optional LibreOffice path when needed."
    ),
    MODE_HTML_TO_MD: (
        "Convert HTML files to Markdown. Pandoc is used when available, with a built-in structured text-and-list fallback when Pandoc is missing."
    ),
    MODE_MERGE_PDFS: "Merge many PDF files into one final PDF.",
    MODE_PRESENTATIONS_TO_PDF: (
        "Convert PPT, PPTX, or ODP files to PDF. You can convert each file separately or merge the resulting PDFs into one final file."
    ),
    MODE_PRESENTATIONS_TO_IMAGES: (
        "Convert PPT, PPTX, or ODP files into slide images. Patch 13 converts each presentation through the same pure Python or optional LibreOffice PDF route, then exports each PDF page as an image."
    ),
    MODE_PDF_TO_PPTX: (
        "Convert each PDF page into a PowerPoint slide. This mode creates one slide per PDF page using page images for layout fidelity."
    ),
    MODE_PDF_TO_HTML: (
        "Export PDFs as standalone HTML documents. This is best-effort and works best on text-heavy PDFs or simple page layouts."
    ),
}

PDF_TOOL_MERGE = "Merge PDFs"
PDF_TOOL_SPLIT_RANGES = "Split PDF by ranges"
PDF_TOOL_SPLIT_EVERY_N = "Split PDF every N pages"
PDF_TOOL_EXTRACT_PAGES = "Extract pages"
PDF_TOOL_REMOVE_PAGES = "Remove pages"
PDF_TOOL_REORDER_PAGES = "Reorder pages"
PDF_TOOL_WATERMARK_TEXT = "Add text watermark"
PDF_TOOL_WATERMARK_IMAGE = "Add image watermark"
PDF_TOOL_TEXT_OVERLAY = "Edit PDF with text overlay"
PDF_TOOL_IMAGE_OVERLAY = "Edit PDF with image overlay"
PDF_TOOL_REDACT_TEXT = "Redact searched text"
PDF_TOOL_REDACT_AREA = "Redact area / region"
PDF_TOOL_EDIT_TEXT = "Edit PDF text (best-effort)"
PDF_TOOL_SIGN_VISIBLE = "Sign PDF (visible)"
PDF_TOOL_EDIT_METADATA = "Edit metadata"
PDF_TOOL_LOCK = "Lock PDF with password"
PDF_TOOL_UNLOCK = "Unlock PDF"
PDF_TOOL_COMPRESS = "Compress PDF"

PDF_TOOL_ORDER = [
    PDF_TOOL_MERGE,
    PDF_TOOL_SPLIT_RANGES,
    PDF_TOOL_SPLIT_EVERY_N,
    PDF_TOOL_EXTRACT_PAGES,
    PDF_TOOL_REMOVE_PAGES,
    PDF_TOOL_REORDER_PAGES,
    PDF_TOOL_WATERMARK_TEXT,
    PDF_TOOL_WATERMARK_IMAGE,
    PDF_TOOL_TEXT_OVERLAY,
    PDF_TOOL_IMAGE_OVERLAY,
    PDF_TOOL_REDACT_TEXT,
    PDF_TOOL_REDACT_AREA,
    PDF_TOOL_EDIT_TEXT,
    PDF_TOOL_SIGN_VISIBLE,
    PDF_TOOL_EDIT_METADATA,
    PDF_TOOL_LOCK,
    PDF_TOOL_UNLOCK,
    PDF_TOOL_COMPRESS,
]

PDF_TOOL_HELP: dict[str, str] = {
    PDF_TOOL_MERGE: "Combine many PDFs into one final file. File order follows the selected input order.",
    PDF_TOOL_SPLIT_RANGES: (
        "Create multiple output PDFs using page groups. Use semicolons between groups, such as "
        "1-3; 4-6; 8; 10-12."
    ),
    PDF_TOOL_SPLIT_EVERY_N: "Split each PDF into sequential parts of N pages each.",
    PDF_TOOL_EXTRACT_PAGES: "Keep only the pages you specify, such as 1-3, 6, 9-last.",
    PDF_TOOL_REMOVE_PAGES: "Remove the specified pages and keep the rest.",
    PDF_TOOL_REORDER_PAGES: (
        "Create a new PDF in the order you specify. This can also duplicate pages. Example: 3,1,2,2,5."
    ),
    PDF_TOOL_WATERMARK_TEXT: "Overlay text on one or many PDFs with configurable size, angle, opacity, and position.",
    PDF_TOOL_WATERMARK_IMAGE: "Overlay a PNG, JPG, or other image as a watermark on one or many PDFs.",
    PDF_TOOL_TEXT_OVERLAY: "Stamp or place editable helper text on selected pages. Leave Pages / ranges blank to target all pages.",
    PDF_TOOL_IMAGE_OVERLAY: "Place an image such as a logo, approval stamp, or screenshot on selected pages.",
    PDF_TOOL_REDACT_TEXT: "Securely search for text and permanently remove matching content from the PDF on selected pages.",
    PDF_TOOL_REDACT_AREA: "Securely redact a rectangular area on one or many pages. Use points or percentages such as 36,72,420,160 or 10%,10%,90%,25%.",
    PDF_TOOL_EDIT_TEXT: "Best-effort text replacement for extractable text only. This uses search plus redaction/replacement and is safest for simple editable PDFs.",
    PDF_TOOL_SIGN_VISIBLE: "Apply a visible signature block using an image, typed signer text, or both. This is a visible sign-off, not certificate signing.",
    PDF_TOOL_EDIT_METADATA: "Set or clear PDF metadata fields such as title, author, subject, and keywords.",
    PDF_TOOL_LOCK: "Create a password-protected PDF using AES-256 encryption. This patch creates a user password and optional owner password for access control.",
    PDF_TOOL_UNLOCK: "Remove PDF password protection by opening the PDF with the correct password and saving an unencrypted copy.",
    PDF_TOOL_COMPRESS: "Reduce PDF size with safe, balanced, or strong lossless compression profiles. For encrypted PDFs, provide the current password.",
}


@dataclass(slots=True)
class BatchConfig:
    mode: str
    files: list[Path]
    output_dir: Path
    merge_to_one_pdf: bool = False
    merged_output_name: str = "converted_output"
    image_format: str = "png"
    image_scale: float = 2.0
    engine_mode: str = ENGINE_AUTO
    soffice_path: str = ""


@dataclass(slots=True)
class PdfToolConfig:
    tool: str
    files: list[Path]
    output_dir: Path
    output_name: str = "merged_pdfs"
    page_spec: str = ""
    every_n_pages: int = 1
    watermark_text: str = "CONFIDENTIAL"
    watermark_image: Path | None = None
    watermark_font_size: int = 42
    watermark_rotation: float = 45.0
    watermark_opacity: float = 0.18
    watermark_position: str = "center"
    watermark_image_scale_percent: int = 40
    metadata_title: str = ""
    metadata_author: str = ""
    metadata_subject: str = ""
    metadata_keywords: str = ""
    metadata_clear_existing: bool = False
    redact_rect: str = ""
    replacement_text: str = ""
    pdf_password: str = ""
    pdf_owner_password: str = ""
    compression_profile: str = "balanced"


def supported_extensions_for_mode(mode: str) -> set[str]:
    return set(MODE_TO_EXTENSIONS.get(mode, set()))


def resolve_soffice_path(soffice_path: str | Path | None = None) -> str | None:
    candidates: list[Path] = []
    if soffice_path:
        base = Path(str(soffice_path)).expanduser()
        if base.is_dir():
            candidates.extend([
                base / "soffice",
                base / "soffice.exe",
                base / "program" / "soffice",
                base / "program" / "soffice.exe",
            ])
        else:
            candidates.append(base)
    for candidate in candidates:
        if candidate.exists():
            return str(candidate.resolve())
    return find_command("libreoffice", "soffice")


def dependency_status(soffice_path: str | Path | None = None) -> dict[str, bool]:
    return {
        "LibreOffice": bool(resolve_soffice_path(soffice_path)),
        "Pandoc": bool(find_command("pandoc")),
        "pdftotext": bool(find_command("pdftotext")),
        "python-pptx": importlib.util.find_spec("pptx") is not None,
        "xlrd": importlib.util.find_spec("xlrd") is not None,
    }


PURE_PYTHON_RICH_EXTS = {".docx", ".xlsx", ".xls", ".csv", ".tsv", ".pptx", ".html", ".htm", ".md", ".markdown", ".mdown"}
PURE_PYTHON_TEXT_FIRST_EXTS = {".txt", ".text", ".log", ".ini", ".cfg", ".conf", ".json", ".xml", ".yaml", ".yml", ".toml", ".rtf", ".odt", ".fodt", ".ods", ".odp"}
PURE_PYTHON_UNSUPPORTED_LEGACY_EXTS = {".doc", ".ppt"}


def describe_pdf_output_route(input_path: Path, engine_mode: str = ENGINE_AUTO, soffice_path: str | Path | None = None) -> str:
    ext = Path(input_path).suffix.lower()
    file_name = Path(input_path).name
    engine = (engine_mode or ENGINE_AUTO).strip().lower() or ENGINE_AUTO
    soffice_ready = bool(resolve_soffice_path(soffice_path))

    if ext in IMAGE_EXTS:
        return f"{file_name}: built-in image raster pipeline -> PDF | fidelity: original image content"
    if ext in PDF_EXTS:
        return f"{file_name}: PDF passthrough copy | fidelity: original PDF preserved"
    if ext in TEXT_EXTS:
        return f"{file_name}: pure Python text renderer | fidelity: text-first"
    if ext in MARKDOWN_EXTS:
        if engine == ENGINE_LIBREOFFICE:
            return f"{file_name}: LibreOffice route via temporary Office conversion | fidelity: higher layout fidelity" if soffice_ready else f"{file_name}: LibreOffice selected but soffice is not configured; pure Python Markdown renderer is recommended"
        if engine == ENGINE_AUTO and soffice_ready:
            return f"{file_name}: pure Python Markdown renderer first, LibreOffice fallback available | fidelity: structured markdown"
        return f"{file_name}: pure Python Markdown renderer | fidelity: structured markdown"
    if ext in WEB_EXTS:
        if engine == ENGINE_LIBREOFFICE:
            return f"{file_name}: LibreOffice HTML route | fidelity: higher layout fidelity" if soffice_ready else f"{file_name}: LibreOffice selected but soffice is not configured; pure Python HTML renderer is recommended"
        if engine == ENGINE_AUTO and soffice_ready:
            return f"{file_name}: pure Python HTML structural renderer first, LibreOffice fallback available | fidelity: structure-aware"
        return f"{file_name}: pure Python HTML structural renderer | fidelity: structure-aware"
    if ext == ".docx":
        if engine == ENGINE_LIBREOFFICE:
            return f"{file_name}: LibreOffice DOCX renderer | fidelity: higher layout fidelity" if soffice_ready else f"{file_name}: LibreOffice selected but soffice is not configured; pure Python DOCX renderer is available"
        if engine == ENGINE_AUTO and soffice_ready:
            return f"{file_name}: pure Python DOCX structure renderer first, LibreOffice fallback available | fidelity: structure-aware"
        return f"{file_name}: pure Python DOCX structure renderer | fidelity: structure-aware"
    if ext in {".xlsx", ".xls", ".csv", ".tsv"}:
        if engine == ENGINE_LIBREOFFICE:
            return f"{file_name}: LibreOffice spreadsheet renderer | fidelity: higher workbook fidelity" if soffice_ready else f"{file_name}: LibreOffice selected but soffice is not configured; pure Python table renderer is available"
        if engine == ENGINE_AUTO and soffice_ready:
            return f"{file_name}: pure Python sheet table renderer first, LibreOffice fallback available | fidelity: table-aware"
        return f"{file_name}: pure Python sheet table renderer | fidelity: table-aware"
    if ext == ".pptx":
        if engine == ENGINE_LIBREOFFICE:
            return f"{file_name}: LibreOffice presentation renderer | fidelity: higher slide fidelity" if soffice_ready else f"{file_name}: LibreOffice selected but soffice is not configured; pure Python slide summary renderer is available"
        if engine == ENGINE_AUTO and soffice_ready:
            return f"{file_name}: pure Python slide summary renderer first, LibreOffice fallback available | fidelity: content-first"
        return f"{file_name}: pure Python slide summary renderer | fidelity: content-first"
    if ext in {".rtf", ".odt", ".fodt", ".ods", ".odp"}:
        if engine == ENGINE_LIBREOFFICE and soffice_ready:
            return f"{file_name}: LibreOffice format renderer | fidelity: higher layout fidelity"
        if engine == ENGINE_AUTO and soffice_ready:
            return f"{file_name}: pure Python text-first fallback first, LibreOffice fallback available | fidelity: text-first"
        return f"{file_name}: pure Python text-first fallback | fidelity: text-first"
    if ext in PURE_PYTHON_UNSUPPORTED_LEGACY_EXTS:
        if soffice_ready:
            return f"{file_name}: legacy Office format; Auto and LibreOffice use soffice for conversion | fidelity: fallback external renderer"
        return f"{file_name}: legacy Office format; configure soffice for conversion because pure Python support is not available"
    if looks_like_text(Path(input_path)):
        return f"{file_name}: detected as text, using pure Python text renderer | fidelity: text-first"
    return f"{file_name}: unsupported for direct PDF output in this mode"


MODE_ROUTE_HELP: dict[str, str] = {
    MODE_PDF_TO_IMAGES: "Uses PyMuPDF page rendering. Engine selection does not affect this mode.",
    MODE_PDF_TO_DOCX: "Uses PDF text extraction with DOCX writing. Best for text-heavy PDFs; engine selection does not affect this mode.",
    MODE_PDF_TO_XLSX: "Uses pdfplumber table extraction with page-text fallback. Engine selection does not affect this mode.",
    MODE_MD_TO_PDF: "Uses the built-in Markdown-to-HTML renderer, then the pure Python PDF renderer.",
    MODE_MD_TO_DOCX: "Uses Pandoc when available, otherwise the built-in Markdown-to-DOCX writer.",
    MODE_MD_TO_HTML: "Uses Pandoc when available, otherwise the built-in Markdown-to-HTML renderer.",
    MODE_HTML_TO_PDF: "Uses the pure Python HTML renderer first, with optional LibreOffice fallback when your engine allows it.",
    MODE_HTML_TO_DOCX: "Uses Pandoc when available, otherwise the built-in structured HTML-to-DOCX renderer, with optional LibreOffice fallback.",
    MODE_HTML_TO_MD: "Uses Pandoc when available, otherwise the built-in HTML-to-Markdown fallback.",
    MODE_PRESENTATIONS_TO_IMAGES: "Converts presentations to PDF using the selected engine path, then renders each resulting page as an image.",
    MODE_PDF_TO_HTML: "Uses PyMuPDF HTML extraction. Engine selection does not affect this mode.",
    MODE_PDF_TO_PPTX: "Uses PyMuPDF page rendering plus python-pptx slide creation. Engine selection does not affect this mode.",
    MODE_MERGE_PDFS: "Uses pypdf merge logic. Engine selection does not affect this mode.",
}


def describe_input_route_for_mode(mode: str, input_path: Path, engine_mode: str = ENGINE_AUTO, soffice_path: str | Path | None = None) -> str:
    if mode in PDF_OUTPUT_MODES:
        return describe_pdf_output_route(input_path, engine_mode=engine_mode, soffice_path=soffice_path)
    return f"{Path(input_path).name}: {MODE_ROUTE_HELP.get(mode, 'Built-in pipeline for the selected mode.')}"


def build_conversion_route_preview(
    mode: str,
    files: Sequence[Path] | None = None,
    *,
    engine_mode: str = ENGINE_AUTO,
    soffice_path: str | Path | None = None,
    limit: int = 4,
) -> str:
    items = list(files or [])
    if not items:
        if mode in PDF_OUTPUT_MODES:
            return (
                "Routing preview: add files to see exact per-file routing. Patch 6 prefers richer pure Python renderers for DOCX, XLSX, PPTX, HTML, Markdown, TXT, CSV, and TSV, then falls back to LibreOffice only when needed."
            )
        return f"Routing preview: {MODE_ROUTE_HELP.get(mode, 'Built-in pipeline for the selected mode.')}"

    lines = [describe_input_route_for_mode(mode, path, engine_mode=engine_mode, soffice_path=soffice_path) for path in items[:limit]]
    if len(items) > limit:
        lines.append(f"... and {len(items) - limit} more file(s)")
    return "Routing preview -> " + " | ".join(lines)


def outputs_pdf(mode: str) -> bool:
    return mode in PDF_OUTPUT_MODES or mode == MODE_MERGE_PDFS


def default_merged_name(mode: str) -> str:
    mapping = {
        MODE_ANY_TO_PDF: "merged_output",
        MODE_IMAGES_TO_PDF: "images_combined",
        MODE_DOCS_TO_PDF: "documents_combined",
        MODE_SHEETS_TO_PDF: "spreadsheets_combined",
        MODE_TEXT_TO_PDF: "text_combined",
        MODE_MD_TO_PDF: "markdown_combined",
        MODE_HTML_TO_PDF: "html_combined",
        MODE_PRESENTATIONS_TO_PDF: "presentations_combined",
        MODE_MERGE_PDFS: "merged_pdfs",
    }
    return mapping.get(mode, "converted_output")


def collect_files_from_folder(folder: Path, allowed_exts: set[str], recursive: bool = False) -> list[Path]:
    folder = Path(folder)
    iterator = folder.rglob("*") if recursive else folder.glob("*")
    files = [p for p in iterator if p.is_file() and p.suffix.lower() in allowed_exts]
    return dedupe_paths(sorted(files, key=lambda p: str(p).lower()))


def dedupe_paths(paths: Iterable[Path]) -> list[Path]:
    result: list[Path] = []
    seen: set[str] = set()
    for path in paths:
        normalized = str(Path(path).expanduser().resolve())
        if normalized not in seen:
            seen.add(normalized)
            result.append(Path(normalized))
    return result


def ensure_directory(path: Path) -> Path:
    path.mkdir(parents=True, exist_ok=True)
    return path


def find_command(*names: str) -> str | None:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def _startupinfo() -> subprocess.STARTUPINFO | None:
    if os.name != "nt":
        return None
    info = subprocess.STARTUPINFO()
    info.dwFlags |= subprocess.STARTF_USESHOWWINDOW
    return info


def run_command(cmd: Sequence[str], cwd: Path | None = None) -> subprocess.CompletedProcess[str]:
    completed = subprocess.run(
        [str(part) for part in cmd],
        cwd=str(cwd) if cwd else None,
        capture_output=True,
        text=True,
        check=False,
        startupinfo=_startupinfo(),
    )
    if completed.returncode != 0:
        details = (completed.stderr or completed.stdout or "Unknown command failure.").strip()
        raise ConversionError(f"Command failed: {' '.join(map(str, cmd))}\n{details}")
    return completed


def unique_path(path: Path) -> Path:
    path = Path(path)
    if not path.exists():
        return path
    stem = path.stem
    suffix = path.suffix
    parent = path.parent
    counter = 1
    while True:
        candidate = parent / f"{stem}_{counter}{suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


def safe_name(text: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", text).strip("._")
    return cleaned or "output"


def filetype_patterns_for_mode(mode: str) -> str:
    exts = sorted(supported_extensions_for_mode(mode))
    if not exts:
        return "*.*"
    return " ".join(f"*{ext}" for ext in exts)




class _BasicHTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []

    def _append(self, text: str) -> None:
        if text:
            self.parts.append(text)

    def handle_starttag(self, tag: str, attrs) -> None:  # type: ignore[override]
        normalized = tag.lower()
        if normalized == "br":
            self._append("\n")
        elif normalized in {"p", "div", "section", "article", "header", "footer", "table", "tr"}:
            self._append("\n")
        elif normalized in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._append("\n")
        elif normalized == "li":
            self._append("\n• ")
        elif normalized in {"ul", "ol"}:
            self._append("\n")
        elif normalized == "td":
            self._append(" ")

    def handle_endtag(self, tag: str) -> None:  # type: ignore[override]
        normalized = tag.lower()
        if normalized in {"p", "div", "section", "article", "header", "footer", "table", "tr", "li"}:
            self._append("\n")
        elif normalized in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._append("\n")

    def handle_data(self, data: str) -> None:  # type: ignore[override]
        collapsed = re.sub(r"[\t\f\v ]+", " ", data)
        if collapsed:
            self._append(collapsed)

    def get_text(self) -> str:
        raw = "".join(self.parts)
        lines = [re.sub(r"[\t ]+", " ", line).strip() for line in raw.splitlines()]
        cleaned: list[str] = []
        blank_pending = False
        for line in lines:
            if line:
                if blank_pending and cleaned:
                    cleaned.append("")
                cleaned.append(line)
                blank_pending = False
            else:
                blank_pending = True
        return "\n".join(cleaned).strip()


def extract_text_from_html_string(html_content: str) -> str:
    parser = _BasicHTMLTextExtractor()
    parser.feed(html_content)
    parser.close()
    return parser.get_text()


def extract_html_title(html_content: str, fallback: str = "") -> str:
    match = re.search(r"<title[^>]*>(.*?)</title>", html_content, flags=re.IGNORECASE | re.DOTALL)
    if not match:
        return fallback
    title = extract_text_from_html_string(match.group(1)).strip()
    return title or fallback


def markdown_inline_to_text(text: str) -> str:
    value = re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", r"\1", text)
    value = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", value)
    value = re.sub(r"`([^`]+)`", r"\1", value)
    value = re.sub(r"\*\*([^*]+)\*\*", r"\1", value)
    value = re.sub(r"__([^_]+)__", r"\1", value)
    value = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", value)
    value = re.sub(r"(?<!_)_([^_]+)_(?!_)", r"\1", value)
    return value.strip()


def _markdown_inline_to_html(text: str) -> str:
    escaped = escape(text)
    escaped = re.sub(
        r"!\[([^\]]*)\]\(([^)]+)\)",
        lambda m: f'<img alt="{escape(m.group(1))}" src="{escape(m.group(2))}" />',
        escaped,
    )
    escaped = re.sub(
        r"\[([^\]]+)\]\(([^)]+)\)",
        lambda m: f'<a href="{escape(m.group(2))}">{m.group(1)}</a>',
        escaped,
    )
    escaped = re.sub(r"`([^`]+)`", r"<code>\1</code>", escaped)
    escaped = re.sub(r"\*\*([^*]+)\*\*", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"__([^_]+)__", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"<em>\1</em>", escaped)
    escaped = re.sub(r"(?<!_)_([^_]+)_(?!_)", r"<em>\1</em>", escaped)
    return escaped


def simple_markdown_to_html_document(markdown_text: str, title: str = "Document") -> str:
    html_parts: list[str] = []
    paragraph_lines: list[str] = []
    list_kind: str | None = None
    in_code = False

    def flush_paragraph() -> None:
        nonlocal paragraph_lines
        if paragraph_lines:
            combined = " ".join(part.strip() for part in paragraph_lines if part.strip()).strip()
            if combined:
                html_parts.append(f"<p>{_markdown_inline_to_html(combined)}</p>")
            paragraph_lines = []

    def close_list() -> None:
        nonlocal list_kind
        if list_kind:
            html_parts.append(f"</{list_kind}>")
            list_kind = None

    for raw_line in markdown_text.splitlines():
        stripped = raw_line.strip()

        if stripped.startswith("```"):
            flush_paragraph()
            close_list()
            if in_code:
                html_parts.append("</code></pre>")
                in_code = False
            else:
                html_parts.append("<pre><code>")
                in_code = True
            continue

        if in_code:
            html_parts.append(escape(raw_line) + "\n")
            continue

        heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading:
            flush_paragraph()
            close_list()
            level = min(len(heading.group(1)), 6)
            html_parts.append(f"<h{level}>{_markdown_inline_to_html(heading.group(2).strip())}</h{level}>")
            continue

        quote = re.match(r"^>\s?(.*)$", stripped)
        if quote:
            flush_paragraph()
            close_list()
            html_parts.append(f"<blockquote>{_markdown_inline_to_html(quote.group(1).strip())}</blockquote>")
            continue

        unordered = re.match(r"^[-*+]\s+(.*)$", stripped)
        if unordered:
            flush_paragraph()
            if list_kind != "ul":
                close_list()
                html_parts.append("<ul>")
                list_kind = "ul"
            html_parts.append(f"<li>{_markdown_inline_to_html(unordered.group(1).strip())}</li>")
            continue

        ordered = re.match(r"^\d+[.)]\s+(.*)$", stripped)
        if ordered:
            flush_paragraph()
            if list_kind != "ol":
                close_list()
                html_parts.append("<ol>")
                list_kind = "ol"
            html_parts.append(f"<li>{_markdown_inline_to_html(ordered.group(1).strip())}</li>")
            continue

        if not stripped:
            flush_paragraph()
            close_list()
            continue

        if list_kind:
            close_list()
        paragraph_lines.append(raw_line)

    flush_paragraph()
    close_list()
    if in_code:
        html_parts.append("</code></pre>")

    body = "\n".join(html_parts)
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <title>{escape(title)}</title>
  <style>
    body {{ font-family: Arial, Helvetica, sans-serif; margin: 2rem auto; max-width: 900px; line-height: 1.55; color: #111827; }}
    h1, h2, h3, h4, h5, h6 {{ margin-top: 1.4em; }}
    pre {{ background: #111827; color: #f9fafb; padding: 1rem; overflow-x: auto; border-radius: 8px; }}
    code {{ background: #f3f4f6; padding: 0.1rem 0.25rem; border-radius: 4px; }}
    blockquote {{ border-left: 4px solid #9ca3af; margin: 1rem 0; padding: 0.5rem 1rem; color: #374151; background: #f9fafb; }}
  </style>
</head>
<body>
{body}
</body>
</html>
"""


def markdown_to_plain_text(markdown_text: str) -> str:
    lines: list[str] = []
    in_code = False
    for raw_line in markdown_text.splitlines():
        stripped = raw_line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            lines.append(raw_line.rstrip())
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading:
            lines.append(markdown_inline_to_text(heading.group(2)))
            lines.append("")
            continue
        quote = re.match(r"^>\s?(.*)$", stripped)
        if quote:
            lines.append(markdown_inline_to_text(quote.group(1)))
            continue
        unordered = re.match(r"^[-*+]\s+(.*)$", stripped)
        if unordered:
            lines.append(f"• {markdown_inline_to_text(unordered.group(1))}")
            continue
        ordered = re.match(r"^(\d+[.)])\s+(.*)$", stripped)
        if ordered:
            lines.append(f"{ordered.group(1)} {markdown_inline_to_text(ordered.group(2))}")
            continue
        if not stripped:
            lines.append("")
            continue
        lines.append(markdown_inline_to_text(raw_line))
    return "\n".join(lines).strip()


def _docx_apply_runs(paragraph, text: str, *, monospace: bool = False) -> None:
    run = paragraph.add_run(text)
    if monospace:
        run.font.name = "Courier New"
        run.font.size = Pt(10)


def save_string_as_docx(text: str, output_docx: Path, title: str | None = None, *, monospace: bool = False) -> Path:
    output_docx = unique_path(output_docx)
    ensure_directory(output_docx.parent)

    document = Document()
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Courier New" if monospace else "Calibri"
    normal_style.font.size = Pt(10 if monospace else 11)

    if title:
        document.add_heading(title, level=0)

    lines = text.splitlines()
    if not lines:
        lines = [""]
    for line in lines:
        paragraph = document.add_paragraph()
        _docx_apply_runs(paragraph, line, monospace=monospace)

    document.save(str(output_docx))
    return output_docx


def simple_markdown_to_docx(input_md: Path, output_docx: Path) -> Path:
    markdown_text = read_text_file(input_md)
    output_docx = unique_path(output_docx)
    ensure_directory(output_docx.parent)

    document = Document()
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Calibri"
    normal_style.font.size = Pt(11)

    paragraph_lines: list[str] = []
    in_code = False

    def flush_paragraph() -> None:
        nonlocal paragraph_lines
        if paragraph_lines:
            paragraph = document.add_paragraph()
            _docx_apply_runs(paragraph, " ".join(paragraph_lines).strip(), monospace=False)
            paragraph_lines = []

    for raw_line in markdown_text.splitlines():
        stripped = raw_line.strip()

        if stripped.startswith("```"):
            flush_paragraph()
            in_code = not in_code
            continue

        if in_code:
            paragraph = document.add_paragraph()
            _docx_apply_runs(paragraph, raw_line.rstrip(), monospace=True)
            continue

        if not stripped:
            flush_paragraph()
            continue

        heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading:
            flush_paragraph()
            level = min(len(heading.group(1)), 6)
            document.add_heading(markdown_inline_to_text(heading.group(2)), level=level)
            continue

        quote = re.match(r"^>\s?(.*)$", stripped)
        if quote:
            flush_paragraph()
            paragraph = document.add_paragraph()
            _docx_apply_runs(paragraph, markdown_inline_to_text(quote.group(1)), monospace=False)
            paragraph.paragraph_format.left_indent = Pt(18)
            continue

        unordered = re.match(r"^[-*+]\s+(.*)$", stripped)
        if unordered:
            flush_paragraph()
            try:
                paragraph = document.add_paragraph(style="List Bullet")
            except Exception:
                paragraph = document.add_paragraph()
            _docx_apply_runs(paragraph, markdown_inline_to_text(unordered.group(1)), monospace=False)
            continue

        ordered = re.match(r"^(\d+[.)])\s+(.*)$", stripped)
        if ordered:
            flush_paragraph()
            try:
                paragraph = document.add_paragraph(style="List Number")
            except Exception:
                paragraph = document.add_paragraph()
            _docx_apply_runs(paragraph, markdown_inline_to_text(ordered.group(2)), monospace=False)
            continue

        paragraph_lines.append(markdown_inline_to_text(raw_line))

    flush_paragraph()
    document.save(str(output_docx))
    return output_docx


def save_string_as_pdf(text: str, output_pdf: Path) -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    pdf = canvas.Canvas(str(output_pdf), pagesize=A4)
    _width, height = A4
    left_margin = 15 * mm
    top_margin = 15 * mm
    bottom_margin = 15 * mm
    y = height - top_margin
    pdf.setFont("Courier", 10)

    max_chars = 95
    line_height = 12

    lines: list[str] = []
    for raw_line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        expanded = raw_line.expandtabs(4)
        wrapped = textwrap.wrap(
            expanded,
            width=max_chars,
            replace_whitespace=False,
            drop_whitespace=False,
            break_long_words=True,
            break_on_hyphens=False,
        )
        lines.extend(wrapped or [""])

    if not lines:
        lines = [""]

    for line in lines:
        if y <= bottom_margin:
            pdf.showPage()
            pdf.setFont("Courier", 10)
            y = height - top_margin
        pdf.drawString(left_margin, y, line)
        y -= line_height

    pdf.save()
    return output_pdf


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def looks_like_text(path: Path) -> bool:
    try:
        sample = path.read_bytes()[:4096]
    except OSError:
        return False
    if not sample:
        return True
    if b"\x00" in sample:
        return False
    text_chars = sum(1 for byte in sample if 32 <= byte <= 126 or byte in b"\n\r\t\f\b")
    return (text_chars / max(len(sample), 1)) > 0.70


def _clamp(value: float, lower: float, upper: float) -> float:
    return max(lower, min(upper, value))


def save_text_as_pdf(input_path: Path, output_pdf: Path) -> Path:
    return save_string_as_pdf(read_text_file(input_path), output_pdf)


def convert_images_to_single_pdf(image_paths: Sequence[Path], output_pdf: Path) -> Path:
    if not image_paths:
        raise ConversionError("No image files were provided.")

    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    opened_images: list[Image.Image] = []
    try:
        for image_path in image_paths:
            image = Image.open(image_path)
            if getattr(image, "is_animated", False):
                image.seek(0)
            if image.mode not in ("RGB", "L"):
                image = image.convert("RGB")
            elif image.mode == "L":
                image = image.convert("RGB")
            opened_images.append(image)

        first, *rest = opened_images
        first.save(str(output_pdf), save_all=True, append_images=rest)
        return output_pdf
    finally:
        for image in opened_images:
            try:
                image.close()
            except Exception:
                pass


def merge_pdfs(pdf_paths: Sequence[Path], output_pdf: Path) -> Path:
    if not pdf_paths:
        raise ConversionError("No PDF files were provided to merge.")

    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    writer = PdfWriter()
    try:
        for pdf_path in pdf_paths:
            if hasattr(writer, "append"):
                writer.append(str(pdf_path))
            else:
                reader = PdfReader(str(pdf_path))
                for page in reader.pages:
                    writer.add_page(page)
        with output_pdf.open("wb") as handle:
            writer.write(handle)
    finally:
        if hasattr(writer, "close"):
            writer.close()
    return output_pdf


def libreoffice_convert_to(input_path: Path, target_ext: str, output_path: Path, soffice_path: str | Path | None = None) -> Path:
    office = resolve_soffice_path(soffice_path)
    if not office:
        raise ConversionError(
            "LibreOffice was not found. Set the soffice path in Settings or add LibreOffice to PATH for Office/PDF conversions."
        )

    output_path = unique_path(output_path)
    ensure_directory(output_path.parent)

    with tempfile.TemporaryDirectory() as out_dir, tempfile.TemporaryDirectory() as profile_dir:
        temp_out = Path(out_dir)
        profile_uri = Path(profile_dir).resolve().as_uri()
        cmd = [
            office,
            f"-env:UserInstallation={profile_uri}",
            "--headless",
            "--convert-to",
            target_ext,
            "--outdir",
            str(temp_out),
            str(input_path),
        ]
        run_command(cmd)

        exact = temp_out / f"{input_path.stem}.{target_ext}"
        if exact.exists():
            shutil.move(str(exact), str(output_path))
            return output_path

        candidates = sorted(
            [p for p in temp_out.iterdir() if p.stem == input_path.stem],
            key=lambda p: p.stat().st_mtime,
            reverse=True,
        )
        if not candidates:
            raise ConversionError(f"LibreOffice did not produce an output file for {input_path.name}.")

        shutil.move(str(candidates[0]), str(output_path))
        return output_path


def markdown_to_docx(input_md: Path, output_docx: Path) -> Path:
    pandoc = find_command("pandoc")
    if pandoc:
        output_docx = unique_path(output_docx)
        ensure_directory(output_docx.parent)
        run_command([pandoc, str(input_md), "-o", str(output_docx)])
        return output_docx

    return simple_markdown_to_docx(input_md, output_docx)


def markdown_to_pdf(input_md: Path, output_pdf: Path) -> Path:
    markdown_text = read_text_file(input_md)
    html_document = simple_markdown_to_html_document(markdown_text, title=input_md.stem)
    return render_html_to_pdf_from_string(html_document, output_pdf, title=input_md.stem)


def markdown_to_html(input_md: Path, output_html: Path) -> Path:
    pandoc = find_command("pandoc")
    output_html = unique_path(output_html)
    ensure_directory(output_html.parent)
    if pandoc:
        run_command([pandoc, str(input_md), "-s", "-o", str(output_html)])
        return output_html

    markdown_text = read_text_file(input_md)
    output_html.write_text(simple_markdown_to_html_document(markdown_text, title=input_md.stem), encoding="utf-8")
    return output_html


def html_to_docx(input_html: Path, output_docx: Path) -> Path:
    pandoc = find_command("pandoc")
    if pandoc:
        output_docx = unique_path(output_docx)
        ensure_directory(output_docx.parent)
        run_command([pandoc, str(input_html), "-o", str(output_docx)])
        return output_docx

    html_content = read_text_file(input_html)
    title = extract_html_title(html_content, fallback=input_html.stem)

    try:
        return render_html_to_docx_from_string(html_content, output_docx, title=title)
    except Exception:
        try:
            return libreoffice_convert_to(input_html, "docx", output_docx, soffice_path=resolve_soffice_path())
        except ConversionError:
            text_value = extract_text_from_html_string(html_content)
            return save_string_as_docx(text_value, output_docx, title=title, monospace=False)



def _html_blocks_to_markdown_lines(html_content: str) -> list[str]:
    parser = _BasicHTMLTextExtractor()
    parser.feed(html_content)
    parser.close()
    text = parser.get_text()
    if not text.strip():
        return ["# Document", "", "_No extractable HTML text was found._"]
    lines: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            if lines and lines[-1] != "":
                lines.append("")
            continue
        if line.startswith("• "):
            lines.append(f"- {line[2:].strip()}")
        else:
            lines.append(line)
    return lines or [text.strip()]


def html_to_markdown(input_html: Path, output_md: Path) -> Path:
    pandoc = find_command("pandoc")
    output_md = unique_path(output_md)
    ensure_directory(output_md.parent)
    if pandoc:
        run_command([pandoc, str(input_html), "-f", "html", "-t", "gfm", "-o", str(output_md)])
        return output_md

    html_content = read_text_file(input_html)
    title = extract_html_title(html_content, fallback=input_html.stem)
    lines = [f"# {title}", ""] if title else []
    lines.extend(_html_blocks_to_markdown_lines(html_content))
    output_md.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")
    return output_md


def html_to_pdf(
    input_html: Path,
    output_pdf: Path,
    *,
    engine_mode: str = ENGINE_AUTO,
    soffice_path: str | Path | None = None,
) -> Path:
    html_content = read_text_file(input_html)
    title = extract_html_title(html_content, fallback=input_html.stem)

    engine = (engine_mode or ENGINE_AUTO).strip().lower()
    if engine == ENGINE_LIBREOFFICE:
        return libreoffice_convert_to(input_html, "pdf", output_pdf, soffice_path=soffice_path)

    pure_python_error: Exception | None = None
    try:
        return render_html_to_pdf_from_string(html_content, output_pdf, title=title or input_html.stem)
    except Exception as exc:
        pure_python_error = exc

    if engine == ENGINE_PURE_PYTHON:
        raise ConversionError(f"Pure Python HTML -> PDF failed for {input_html.name}: {pure_python_error}") from pure_python_error

    try:
        return libreoffice_convert_to(input_html, "pdf", output_pdf, soffice_path=soffice_path)
    except Exception as libreoffice_error:
        if pure_python_error is not None:
            raise ConversionError(
                f"Built-in HTML rendering and LibreOffice fallback both failed for {input_html.name}.\n\nBuilt-in error: {pure_python_error}\n\nLibreOffice error: {libreoffice_error}"
            ) from libreoffice_error
        raise


def presentation_to_images(
    input_presentation: Path,
    output_dir: Path,
    *,
    image_format: str = "png",
    image_scale: float = 2.0,
    engine_mode: str = ENGINE_AUTO,
    soffice_path: str | Path | None = None,
) -> list[Path]:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_pdf = Path(temp_dir) / f"{safe_name(input_presentation.stem)}.pdf"
        office_like_to_pdf(input_presentation, temp_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
        presentation_output_dir = ensure_directory(Path(output_dir) / safe_name(input_presentation.stem))
        return pdf_to_images(temp_pdf, presentation_output_dir, image_format=image_format, image_scale=image_scale)


def pdf_to_html(input_pdf: Path, output_html: Path) -> Path:
    output_html = unique_path(output_html)
    ensure_directory(output_html.parent)

    sections: list[str] = []
    with fitz.open(str(input_pdf)) as document:
        if document.page_count < 1:
            raise ConversionError("This PDF has no pages.")
        for page_number, page in enumerate(document, start=1):
            section_html = page.get_text("html") or f"<p>No extractable content found on page {page_number}.</p>"
            sections.append(
                f'<section class="page-shell" data-page="{page_number}">\n'
                f'  <div class="page-label">Page {page_number}</div>\n'
                f'  <div class="page-fragment">{section_html}</div>\n'
                f'</section>'
            )

    html_document = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <title>{escape(input_pdf.stem)}</title>
  <style>
    body {{ margin: 0; padding: 2rem; background: #0f172a; color: #e5e7eb; font-family: Arial, Helvetica, sans-serif; }}
    .page-shell {{ background: #ffffff; color: #111827; margin: 0 auto 1.5rem auto; padding: 1rem; border-radius: 14px; max-width: 1100px; box-shadow: 0 18px 50px rgba(15, 23, 42, 0.35); overflow: auto; }}
    .page-label {{ font-size: 0.95rem; font-weight: 700; color: #475569; margin-bottom: 0.75rem; border-bottom: 1px solid #cbd5e1; padding-bottom: 0.5rem; }}
    .page-fragment {{ overflow: auto; }}
    .page-fragment div[id^="page"] {{ margin: 0 auto; }}
  </style>
</head>
<body>
{''.join(sections)}
</body>
</html>
"""
    output_html.write_text(html_document, encoding="utf-8")
    return output_html


def _choose_blank_slide_layout(presentation) -> object:
    for layout in presentation.slide_layouts:
        if getattr(layout, 'name', '').strip().lower() == 'blank':
            return layout
    if len(presentation.slide_layouts) > 6:
        return presentation.slide_layouts[6]
    return presentation.slide_layouts[-1]


def _fit_size_within_box(source_width: int, source_height: int, box_width: int, box_height: int) -> tuple[int, int]:
    if source_width <= 0 or source_height <= 0:
        return box_width, box_height
    ratio = min(box_width / source_width, box_height / source_height)
    return max(int(source_width * ratio), 1), max(int(source_height * ratio), 1)


def pdf_to_pptx(input_pdf: Path, output_pptx: Path, image_scale: float = 2.0) -> Path:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ConversionError(
            "python-pptx is required for PDF -> PPTX conversion. Install the packages from requirements.txt and try again."
        ) from exc

    output_pptx = unique_path(output_pptx)
    ensure_directory(output_pptx.parent)

    with fitz.open(str(input_pdf)) as document:
        if document.page_count < 1:
            raise ConversionError("This PDF has no pages.")

        presentation = Presentation()
        first_rect = document.load_page(0).rect
        presentation.slide_width = max(int(first_rect.width * 12700), 1000000)
        presentation.slide_height = max(int(first_rect.height * 12700), 1000000)
        blank_layout = _choose_blank_slide_layout(presentation)
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)

        for page in document:
            slide = presentation.slides.add_slide(blank_layout)
            pix = page.get_pixmap(matrix=fitz.Matrix(image_scale, image_scale), alpha=False)
            width, height = _fit_size_within_box(pix.width, pix.height, slide_width, slide_height)
            left = int((slide_width - width) / 2)
            top = int((slide_height - height) / 2)
            image_stream = BytesIO(pix.tobytes("png"))
            image_stream.seek(0)
            slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

        presentation.save(str(output_pptx))

    return output_pptx


def pdf_to_images(input_pdf: Path, output_dir: Path, image_format: str = "png", image_scale: float = 2.0) -> list[Path]:
    image_format = image_format.lower().strip()
    if image_format not in {"png", "jpg", "jpeg"}:
        raise ConversionError("Image format must be png or jpg.")

    output_dir = ensure_directory(output_dir)
    created: list[Path] = []

    with fitz.open(str(input_pdf)) as document:
        for page_number, page in enumerate(document, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(image_scale, image_scale), alpha=False)
            mode = "RGB"
            image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
            extension = "jpg" if image_format in {"jpg", "jpeg"} else "png"
            output_path = output_dir / f"{safe_name(input_pdf.stem)}_page_{page_number}.{extension}"
            if extension == "jpg":
                image.save(output_path, quality=95)
            else:
                image.save(output_path)
            created.append(output_path)
            image.close()

    return created


def extract_pdf_pages_as_text(input_pdf: Path) -> list[str]:
    pdftotext = find_command("pdftotext")
    if pdftotext:
        try:
            completed = run_command([pdftotext, "-layout", "-enc", "UTF-8", str(input_pdf), "-"])
            text = completed.stdout or ""
            pages = [page.rstrip() for page in text.split("\f")]
            pages = [page for page in pages if page.strip()]
            if pages:
                return pages
        except ConversionError:
            pass

    pages: list[str] = []
    with fitz.open(str(input_pdf)) as document:
        for page in document:
            pages.append(page.get_text("text", sort=True).rstrip())
    return pages


def pdf_to_docx(input_pdf: Path, output_docx: Path) -> Path:
    pages = extract_pdf_pages_as_text(input_pdf)
    output_docx = unique_path(output_docx)
    ensure_directory(output_docx.parent)

    document = Document()
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Courier New"
    normal_style.font.size = Pt(10)

    for page_index, page_text in enumerate(pages):
        lines = page_text.splitlines() or [""]
        for line in lines:
            paragraph = document.add_paragraph()
            run = paragraph.add_run(line)
            run.font.name = "Courier New"
            run.font.size = Pt(10)
        if page_index < len(pages) - 1:
            document.add_page_break()

    document.save(str(output_docx))
    return output_docx


def _safe_sheet_title(base: str, existing: set[str]) -> str:
    cleaned = re.sub(r"[\\/*?:\[\]]", "_", base).strip() or "Sheet"
    cleaned = cleaned[:31]
    candidate = cleaned
    counter = 1
    while candidate in existing:
        suffix = f"_{counter}"
        candidate = f"{cleaned[: 31 - len(suffix)]}{suffix}"
        counter += 1
    existing.add(candidate)
    return candidate


def _autosize_columns(worksheet) -> None:
    for column_cells in worksheet.columns:
        max_length = 0
        column_letter = column_cells[0].column_letter
        for cell in column_cells:
            if cell.value is None:
                continue
            max_length = max(max_length, len(str(cell.value)))
        worksheet.column_dimensions[column_letter].width = min(max(max_length + 2, 10), 40)


def pdf_to_xlsx(input_pdf: Path, output_xlsx: Path) -> Path:
    output_xlsx = unique_path(output_xlsx)
    ensure_directory(output_xlsx.parent)

    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)
    used_titles: set[str] = set()
    created_any_sheet = False

    with pdfplumber.open(str(input_pdf)) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables() or []

            if tables:
                for table_index, table in enumerate(tables, start=1):
                    title = _safe_sheet_title(f"Page_{page_index}_Table_{table_index}", used_titles)
                    ws = workbook.create_sheet(title)
                    for row in table:
                        ws.append([(cell or "") for cell in row])
                    _autosize_columns(ws)
                    created_any_sheet = True
            else:
                text = page.extract_text(layout=True) or page.extract_text() or ""
                if text.strip():
                    title = _safe_sheet_title(f"Page_{page_index}_Text", used_titles)
                    ws = workbook.create_sheet(title)
                    for raw_line in text.splitlines():
                        if not raw_line.strip():
                            ws.append([""])
                            continue
                        parts = [part.strip() for part in re.split(r"\s{2,}|\t", raw_line.rstrip())]
                        ws.append(parts or [raw_line.strip()])
                    _autosize_columns(ws)
                    created_any_sheet = True

    if not created_any_sheet:
        ws = workbook.create_sheet(_safe_sheet_title("Extracted_Text", used_titles))
        ws["A1"] = "No extractable tables or text were found in this PDF."
        _autosize_columns(ws)

    workbook.save(str(output_xlsx))
    return output_xlsx


def _normalise_text_block(text: str) -> str:
    return text.replace("\xa0", " ").replace("\u200b", "").replace("\r\n", "\n").replace("\r", "\n")


def _strip_rtf_to_text(rtf_text: str) -> str:
    text = re.sub(r"\\'([0-9a-fA-F]{2})", lambda match: bytes.fromhex(match.group(1)).decode("latin-1", errors="ignore"), rtf_text)
    text = text.replace("\\par", "\n").replace("\\line", "\n").replace("\\tab", "\t")
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    return _normalise_text_block(text)


def _extract_text_from_xml_string(xml_text: str) -> str:
    try:
        root = ET.fromstring(xml_text)
    except ET.ParseError:
        return re.sub(r"<[^>]+>", " ", xml_text)

    lines: list[str] = []

    def walk(node) -> None:
        tag = node.tag.split("}")[-1] if isinstance(node.tag, str) else ""
        text_value = (node.text or "").strip()
        if text_value:
            lines.append(text_value)
        for child in list(node):
            walk(child)
        tail_value = (node.tail or "").strip()
        if tail_value:
            lines.append(tail_value)
        if tag in {"p", "h", "line-break", "tab", "table-row", "section", "table"}:
            lines.append("\n")

    walk(root)
    merged: list[str] = []
    current: list[str] = []
    for item in lines:
        if item == "\n":
            line = " ".join(current).strip()
            if line:
                merged.append(line)
            current = []
            continue
        current.append(item)
    if current:
        line = " ".join(current).strip()
        if line:
            merged.append(line)
    return _normalise_text_block("\n".join(merged))


def _extract_text_from_odf_document(input_path: Path) -> str:
    ext = input_path.suffix.lower()
    if ext == ".fodt":
        return _extract_text_from_xml_string(read_text_file(input_path))
    try:
        with zipfile.ZipFile(input_path, "r") as archive:
            xml_bytes = archive.read("content.xml")
    except Exception as exc:
        raise ConversionError(f"Could not read ODF content from {input_path.name}: {exc}") from exc
    return _extract_text_from_xml_string(xml_bytes.decode("utf-8", errors="replace"))


def document_to_plain_text(input_path: Path) -> str:
    ext = input_path.suffix.lower()
    heading = f"Document: {input_path.name}\n{'=' * (len(input_path.name) + 10)}\n"

    if ext == ".docx":
        document = Document(str(input_path))
        blocks: list[str] = [heading]
        table_counter = 0

        for paragraph in document.paragraphs:
            text_value = paragraph.text.strip()
            if not text_value:
                continue
            style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
            if style_name.lower().startswith("heading"):
                blocks.append(text_value)
                blocks.append("-" * len(text_value))
            else:
                blocks.append(text_value)

        for table in document.tables:
            table_counter += 1
            blocks.append(f"\nTable {table_counter}:")
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                blocks.append("\t".join(values).rstrip("\t"))

        return _normalise_text_block("\n".join(blocks).strip() + "\n")

    if ext == ".rtf":
        return heading + _strip_rtf_to_text(read_text_file(input_path))

    if ext in {".odt", ".fodt"}:
        return heading + _extract_text_from_odf_document(input_path)

    if ext == ".doc":
        raise ConversionError(
            "Pure Python conversion for .doc is not available in this patch. Use Auto or LibreOffice mode for legacy Word files."
        )

    raise ConversionError(f"Unsupported document type for pure Python PDF conversion: {input_path.name}")


def _rows_to_text(rows: Iterable[Sequence[object]], *, heading: str) -> str:
    lines: list[str] = [heading]
    for row in rows:
        values = ["" if value is None else str(value) for value in row]
        if any(value.strip() for value in values):
            lines.append("\t".join(values).rstrip("\t"))
    return _normalise_text_block("\n".join(lines).strip() + "\n")


def spreadsheet_to_plain_text(input_path: Path) -> str:
    ext = input_path.suffix.lower()
    heading = f"Spreadsheet: {input_path.name}\n{'=' * (len(input_path.name) + 13)}\n"

    if ext in {".csv", ".tsv"}:
        delimiter = "\t" if ext == ".tsv" else ","
        with input_path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            reader = csv.reader(handle, delimiter=delimiter)
            return _rows_to_text(reader, heading=heading)

    if ext == ".xlsx":
        workbook = load_workbook(filename=str(input_path), data_only=True, read_only=True)
        try:
            blocks: list[str] = [heading]
            for sheet in workbook.worksheets:
                blocks.append(f"\nSheet: {sheet.title}")
                blocks.append("-" * (len(sheet.title) + 7))
                for row in sheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    if any(value.strip() for value in values):
                        blocks.append("\t".join(values).rstrip("\t"))
            return _normalise_text_block("\n".join(blocks).strip() + "\n")
        finally:
            workbook.close()

    if ext == ".xls":
        try:
            import xlrd  # type: ignore
        except Exception as exc:
            raise ConversionError(
                "Pure Python .xls conversion requires the xlrd package. Install requirements.txt or use Auto / LibreOffice mode."
            ) from exc
        book = xlrd.open_workbook(str(input_path), on_demand=True)
        try:
            blocks: list[str] = [heading]
            for sheet in book.sheets():
                blocks.append(f"\nSheet: {sheet.name}")
                blocks.append("-" * (len(sheet.name) + 7))
                for row_index in range(sheet.nrows):
                    values = [sheet.cell_value(row_index, col_index) for col_index in range(sheet.ncols)]
                    values = ["" if value is None else str(value) for value in values]
                    if any(value.strip() for value in values):
                        blocks.append("\t".join(values).rstrip("\t"))
            return _normalise_text_block("\n".join(blocks).strip() + "\n")
        finally:
            book.release_resources()

    if ext == ".ods":
        return heading + _extract_text_from_odf_document(input_path)

    raise ConversionError(f"Unsupported spreadsheet type for pure Python PDF conversion: {input_path.name}")


def presentation_to_plain_text(input_path: Path) -> str:
    ext = input_path.suffix.lower()
    heading = f"Presentation: {input_path.name}\n{'=' * (len(input_path.name) + 14)}\n"

    if ext == ".pptx":
        try:
            from pptx import Presentation
        except Exception as exc:
            raise ConversionError(
                "python-pptx is required for pure Python PPTX conversion. Install requirements.txt and try again."
            ) from exc
        presentation = Presentation(str(input_path))
        blocks: list[str] = [heading]
        for slide_index, slide in enumerate(presentation.slides, start=1):
            blocks.append(f"\nSlide {slide_index}")
            blocks.append("-" * (len(str(slide_index)) + 6))
            added_any_text = False
            for shape in slide.shapes:
                text_value = getattr(shape, "text", "")
                if text_value and text_value.strip():
                    blocks.append(text_value.strip())
                    added_any_text = True
            if not added_any_text:
                blocks.append("[No extractable text on this slide]")
        return _normalise_text_block("\n".join(blocks).strip() + "\n")

    if ext == ".odp":
        return heading + _extract_text_from_odf_document(input_path)

    if ext == ".ppt":
        raise ConversionError(
            "Pure Python conversion for .ppt is not available in this patch. Use Auto or LibreOffice mode for legacy PowerPoint files."
        )

    raise ConversionError(f"Unsupported presentation type for pure Python PDF conversion: {input_path.name}")


def pure_python_file_to_pdf(input_path: Path, output_pdf: Path) -> Path:
    ext = input_path.suffix.lower()
    if ext == ".docx":
        return render_docx_to_pdf(input_path, output_pdf)
    if ext in {".xlsx", ".xls", ".csv", ".tsv"}:
        return render_spreadsheet_to_pdf(input_path, output_pdf)
    if ext == ".pptx":
        return render_presentation_to_pdf(input_path, output_pdf)
    if ext in MARKDOWN_EXTS:
        markdown_text = read_text_file(input_path)
        html_document = simple_markdown_to_html_document(markdown_text, title=input_path.stem)
        return render_html_to_pdf_from_string(html_document, output_pdf, title=input_path.stem)
    if ext in WEB_EXTS:
        html_content = read_text_file(input_path)
        title = extract_html_title(html_content, fallback=input_path.stem)
        return render_html_to_pdf_from_string(html_content, output_pdf, title=title or input_path.stem)
    if ext in DOC_EXTS:
        return save_string_as_pdf(document_to_plain_text(input_path), output_pdf)
    if ext in SHEET_EXTS:
        return save_string_as_pdf(spreadsheet_to_plain_text(input_path), output_pdf)
    if ext in PRESENTATION_EXTS:
        return save_string_as_pdf(presentation_to_plain_text(input_path), output_pdf)
    if ext in TEXT_EXTS or looks_like_text(input_path):
        return save_text_as_pdf(input_path, output_pdf)
    raise ConversionError(f"Pure Python PDF conversion is not available for {input_path.name}.")


def office_like_to_pdf(input_path: Path, output_pdf: Path, engine_mode: str = ENGINE_AUTO, soffice_path: str | Path | None = None) -> Path:
    ext = input_path.suffix.lower()
    office_friendly = ext in DOC_EXTS | SHEET_EXTS | PRESENTATION_EXTS | WEB_EXTS | TEXT_EXTS | MARKDOWN_EXTS
    if not office_friendly:
        raise ConversionError(f"Unsupported Office-like file for PDF conversion: {input_path.name}")

    engine = (engine_mode or ENGINE_AUTO).strip().lower()

    if engine == ENGINE_PURE_PYTHON:
        return pure_python_file_to_pdf(input_path, output_pdf)

    if engine == ENGINE_LIBREOFFICE:
        return libreoffice_convert_to(input_path, "pdf", output_pdf, soffice_path=soffice_path)

    pure_python_error: Exception | None = None
    try:
        return pure_python_file_to_pdf(input_path, output_pdf)
    except Exception as exc:
        pure_python_error = exc

    try:
        return libreoffice_convert_to(input_path, "pdf", output_pdf, soffice_path=soffice_path)
    except Exception as libreoffice_error:
        if pure_python_error is not None:
            raise ConversionError(
                f"Built-in conversion and LibreOffice fallback both failed for {input_path.name}.\n\nBuilt-in error: {pure_python_error}\n\nLibreOffice error: {libreoffice_error}"
            ) from libreoffice_error
        raise


def copy_as_pdf(input_pdf: Path, output_pdf: Path) -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)
    shutil.copy2(input_pdf, output_pdf)
    return output_pdf


def file_to_pdf(input_path: Path, output_pdf: Path, engine_mode: str = ENGINE_AUTO, soffice_path: str | Path | None = None) -> Path:
    ext = input_path.suffix.lower()
    if ext in IMAGE_EXTS:
        return convert_images_to_single_pdf([input_path], output_pdf)
    if ext in PDF_EXTS:
        return copy_as_pdf(input_path, output_pdf)
    if ext in DOC_EXTS | SHEET_EXTS | TEXT_EXTS | WEB_EXTS | PRESENTATION_EXTS | MARKDOWN_EXTS:
        return office_like_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
    if looks_like_text(input_path):
        return save_text_as_pdf(input_path, output_pdf)
    raise ConversionError(f"Unsupported input type for PDF conversion: {input_path.name}")


def convert_input_to_pdf_for_mode(
    mode: str,
    input_path: Path,
    output_pdf: Path,
    *,
    engine_mode: str = ENGINE_AUTO,
    soffice_path: str | Path | None = None,
) -> Path:
    ext = input_path.suffix.lower()

    if mode == MODE_ANY_TO_PDF:
        return file_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
    if mode == MODE_IMAGES_TO_PDF:
        if ext not in IMAGE_EXTS:
            raise ConversionError(f"{input_path.name} is not a supported image file.")
        return convert_images_to_single_pdf([input_path], output_pdf)
    if mode == MODE_DOCS_TO_PDF:
        if ext not in DOC_EXTS:
            raise ConversionError(f"{input_path.name} is not a supported document file.")
        return office_like_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
    if mode == MODE_SHEETS_TO_PDF:
        if ext not in SHEET_EXTS:
            raise ConversionError(f"{input_path.name} is not a supported spreadsheet file.")
        return office_like_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
    if mode == MODE_PRESENTATIONS_TO_PDF:
        if ext not in PRESENTATION_EXTS:
            raise ConversionError(f"{input_path.name} is not a supported presentation file.")
        return office_like_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
    if mode == MODE_TEXT_TO_PDF:
        if ext in MARKDOWN_EXTS | TEXT_EXTS | WEB_EXTS:
            return office_like_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)
        if looks_like_text(input_path):
            return save_text_as_pdf(input_path, output_pdf)
        raise ConversionError(f"{input_path.name} is not a supported text-like file.")
    if mode == MODE_MD_TO_PDF:
        if ext not in MARKDOWN_EXTS:
            raise ConversionError(f"{input_path.name} is not a supported Markdown file.")
        return markdown_to_pdf(input_path, output_pdf)
    if mode == MODE_HTML_TO_PDF:
        if ext not in WEB_EXTS:
            raise ConversionError(f"{input_path.name} is not a supported HTML file.")
        return html_to_pdf(input_path, output_pdf, engine_mode=engine_mode, soffice_path=soffice_path)

    raise ConversionError(f"Mode does not output PDF: {mode}")


def _log(log: LogFn | None, message: str) -> None:
    if log:
        log(message)


def _progress(progress: ProgressFn | None, current: int, total: int) -> None:
    if progress:
        progress(current, total)


def _parse_page_number(raw: str, total_pages: int) -> int:
    token = raw.strip().lower()
    if token in {"end", "last"}:
        return total_pages
    if not token.isdigit():
        raise ConversionError(f"Invalid page reference: {raw}")
    page_number = int(token)
    if not 1 <= page_number <= total_pages:
        raise ConversionError(f"Page {page_number} is outside the valid range 1-{total_pages}.")
    return page_number


def parse_page_spec(spec: str, total_pages: int, allow_duplicates: bool = False) -> list[int]:
    if total_pages < 1:
        raise ConversionError("This PDF has no pages.")

    raw_parts = [part.strip() for part in spec.split(",") if part.strip()]
    if not raw_parts:
        raise ConversionError("A page specification is required.")

    pages: list[int] = []
    seen: set[int] = set()

    for raw_part in raw_parts:
        token = re.sub(r"\s+", "", raw_part)
        if "-" in token:
            start_token, end_token = token.split("-", 1)
            if not start_token or not end_token:
                raise ConversionError(f"Invalid page range: {raw_part}")
            start = _parse_page_number(start_token, total_pages)
            end = _parse_page_number(end_token, total_pages)
            if start > end:
                raise ConversionError(f"Page range must be ascending: {raw_part}")
            sequence = range(start - 1, end)
        else:
            sequence = [(_parse_page_number(token, total_pages) - 1)]

        for page_index in sequence:
            if allow_duplicates or page_index not in seen:
                pages.append(page_index)
                seen.add(page_index)

    if not pages:
        raise ConversionError("No valid pages were selected.")
    return pages


def parse_page_groups(spec: str, total_pages: int) -> list[list[int]]:
    groups_raw = [part.strip() for part in re.split(r"[;|\n]+", spec) if part.strip()]
    if not groups_raw:
        raise ConversionError(
            "A split range specification is required. Example: 1-3; 4-6; 7-9"
        )
    return [parse_page_spec(group, total_pages, allow_duplicates=False) for group in groups_raw]


def resolve_target_pages(
    page_spec: str,
    total_pages: int,
    *,
    default: str = "all",
    allow_duplicates: bool = False,
) -> list[int]:
    cleaned = page_spec.strip()
    if cleaned:
        return parse_page_spec(cleaned, total_pages, allow_duplicates=allow_duplicates)
    if default == "last":
        return [total_pages - 1]
    return list(range(total_pages))


def page_indices_to_label(indices: Sequence[int]) -> str:
    if not indices:
        return "pages"

    labels: list[str] = []
    numbers = [index + 1 for index in indices]
    start = prev = numbers[0]

    for current in numbers[1:]:
        if current == prev + 1:
            prev = current
            continue
        labels.append(f"{start}-{prev}" if start != prev else f"{start}")
        start = prev = current

    labels.append(f"{start}-{prev}" if start != prev else f"{start}")
    return "p" + "_".join(labels)


def _write_pdf_from_reader(reader: PdfReader, page_indices: Sequence[int], output_pdf: Path) -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    writer = PdfWriter()
    try:
        for page_index in page_indices:
            writer.add_page(reader.pages[page_index])
        with output_pdf.open("wb") as handle:
            writer.write(handle)
    finally:
        if hasattr(writer, "close"):
            writer.close()
    return output_pdf


def split_pdf_by_ranges(input_pdf: Path, groups_spec: str, output_dir: Path) -> list[Path]:
    reader = PdfReader(str(input_pdf))
    groups = parse_page_groups(groups_spec, len(reader.pages))
    outputs: list[Path] = []

    for group_index, page_group in enumerate(groups, start=1):
        label = page_indices_to_label(page_group)
        output_pdf = output_dir / f"{safe_name(input_pdf.stem)}_split_{group_index:03d}_{label}.pdf"
        outputs.append(_write_pdf_from_reader(reader, page_group, output_pdf))

    return outputs


def split_pdf_every_n(input_pdf: Path, pages_per_split: int, output_dir: Path) -> list[Path]:
    if pages_per_split < 1:
        raise ConversionError("Pages per split must be at least 1.")

    reader = PdfReader(str(input_pdf))
    total_pages = len(reader.pages)
    outputs: list[Path] = []

    for start_index in range(0, total_pages, pages_per_split):
        page_group = list(range(start_index, min(start_index + pages_per_split, total_pages)))
        label = page_indices_to_label(page_group)
        part_number = (start_index // pages_per_split) + 1
        output_pdf = output_dir / f"{safe_name(input_pdf.stem)}_part_{part_number:03d}_{label}.pdf"
        outputs.append(_write_pdf_from_reader(reader, page_group, output_pdf))

    return outputs


def extract_pdf_pages(input_pdf: Path, page_spec: str, output_pdf: Path) -> Path:
    reader = PdfReader(str(input_pdf))
    page_indices = parse_page_spec(page_spec, len(reader.pages), allow_duplicates=True)
    return _write_pdf_from_reader(reader, page_indices, output_pdf)


def remove_pdf_pages(input_pdf: Path, page_spec: str, output_pdf: Path) -> Path:
    reader = PdfReader(str(input_pdf))
    total_pages = len(reader.pages)
    to_remove = set(parse_page_spec(page_spec, total_pages, allow_duplicates=False))
    remaining = [index for index in range(total_pages) if index not in to_remove]
    if not remaining:
        raise ConversionError("Removing those pages would result in an empty PDF.")
    return _write_pdf_from_reader(reader, remaining, output_pdf)


def reorder_pdf_pages(input_pdf: Path, page_spec: str, output_pdf: Path) -> Path:
    reader = PdfReader(str(input_pdf))
    page_indices = parse_page_spec(page_spec, len(reader.pages), allow_duplicates=True)
    return _write_pdf_from_reader(reader, page_indices, output_pdf)


def _anchor_rect(page_width: float, page_height: float, obj_width: float, obj_height: float, position: str) -> tuple[float, float]:
    position = position.strip().lower() or "center"
    margin = min(page_width, page_height) * 0.06

    if position == "top-left":
        return margin, page_height - margin - obj_height
    if position == "top-right":
        return page_width - margin - obj_width, page_height - margin - obj_height
    if position == "bottom-left":
        return margin, margin
    if position == "bottom-right":
        return page_width - margin - obj_width, margin
    return (page_width - obj_width) / 2, (page_height - obj_height) / 2


def _anchor_point(page_width: float, page_height: float, position: str) -> tuple[float, float]:
    x, y = _anchor_rect(page_width, page_height, 0, 0, position)
    if position.strip().lower() == "center":
        return page_width / 2, page_height / 2
    return x, y


def _copy_pdf_metadata(reader: PdfReader, writer: PdfWriter) -> None:
    try:
        metadata = {
            str(key): str(value)
            for key, value in dict(reader.metadata or {}).items()
            if key is not None and value is not None
        }
        if metadata:
            writer.add_metadata(metadata)
    except Exception:
        pass


def _build_text_watermark_page(
    page_width: float,
    page_height: float,
    text: str,
    font_size: int,
    rotation: float,
    opacity: float,
    position: str,
    color: tuple[float, float, float] = (0.42, 0.42, 0.42),
    font_name: str = "Helvetica-Bold",
):
    buffer = BytesIO()
    overlay = canvas.Canvas(buffer, pagesize=(page_width, page_height))
    overlay.saveState()
    if hasattr(overlay, "setFillAlpha"):
        try:
            overlay.setFillAlpha(opacity)
        except Exception:
            pass
    overlay.setFillColorRGB(*color)
    overlay.setFont(font_name, max(int(font_size), 8))
    x, y = _anchor_point(page_width, page_height, position)
    overlay.translate(x, y)
    overlay.rotate(rotation)
    overlay.drawCentredString(0, 0, text)
    overlay.restoreState()
    overlay.save()
    buffer.seek(0)
    return PdfReader(buffer).pages[0]


def _build_image_watermark_page(
    page_width: float,
    page_height: float,
    image_path: Path,
    opacity: float,
    scale_percent: int,
    position: str,
):
    buffer = BytesIO()
    overlay = canvas.Canvas(buffer, pagesize=(page_width, page_height))

    with Image.open(image_path) as source_image:
        image = source_image.convert("RGBA")
        opacity = _clamp(opacity, 0.01, 1.0)
        if opacity < 1.0:
            alpha = image.getchannel("A")
            alpha = alpha.point(lambda pixel: int(pixel * opacity))
            image.putalpha(alpha)

        image_reader = ImageReader(image)
        img_width, img_height = image.size

        scale_percent = max(int(scale_percent), 1)
        target_width = page_width * (scale_percent / 100.0)
        target_height = target_width * (img_height / max(img_width, 1))

        max_height = page_height * 0.85
        if target_height > max_height:
            shrink = max_height / target_height
            target_width *= shrink
            target_height *= shrink

        x, y = _anchor_rect(page_width, page_height, target_width, target_height, position)
        overlay.drawImage(
            image_reader,
            x,
            y,
            width=target_width,
            height=target_height,
            mask="auto",
            preserveAspectRatio=True,
        )

    overlay.save()
    buffer.seek(0)
    return PdfReader(buffer).pages[0]


def _apply_overlay_to_pages(
    reader: PdfReader,
    output_pdf: Path,
    target_pages: Sequence[int],
    overlay_builder,
) -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    writer = PdfWriter()
    target_set = set(target_pages)
    try:
        _copy_pdf_metadata(reader, writer)
        for page_index, page in enumerate(reader.pages):
            if page_index in target_set:
                overlay_page = overlay_builder(page_index, page)
                page.merge_page(overlay_page)
            writer.add_page(page)
        with output_pdf.open("wb") as handle:
            writer.write(handle)
    finally:
        if hasattr(writer, "close"):
            writer.close()
    return output_pdf


def apply_text_watermark(
    input_pdf: Path,
    output_pdf: Path,
    text: str,
    font_size: int = 42,
    rotation: float = 45.0,
    opacity: float = 0.18,
    position: str = "center",
    page_spec: str = "",
) -> Path:
    if not text.strip():
        raise ConversionError("Watermark text cannot be empty.")

    reader = PdfReader(str(input_pdf))
    opacity = _clamp(float(opacity), 0.01, 1.0)
    target_pages = resolve_target_pages(page_spec, len(reader.pages), default="all")

    return _apply_overlay_to_pages(
        reader,
        output_pdf,
        target_pages,
        lambda _page_index, page: _build_text_watermark_page(
            float(page.mediabox.width),
            float(page.mediabox.height),
            text=text,
            font_size=font_size,
            rotation=rotation,
            opacity=opacity,
            position=position,
        ),
    )


def apply_image_watermark(
    input_pdf: Path,
    output_pdf: Path,
    image_path: Path,
    opacity: float = 0.18,
    scale_percent: int = 40,
    position: str = "center",
    page_spec: str = "",
) -> Path:
    image_path = Path(image_path)
    if not image_path.exists():
        raise ConversionError(f"Watermark image was not found: {image_path}")

    reader = PdfReader(str(input_pdf))
    opacity = _clamp(float(opacity), 0.01, 1.0)
    scale_percent = max(int(scale_percent), 1)
    target_pages = resolve_target_pages(page_spec, len(reader.pages), default="all")

    return _apply_overlay_to_pages(
        reader,
        output_pdf,
        target_pages,
        lambda _page_index, page: _build_image_watermark_page(
            float(page.mediabox.width),
            float(page.mediabox.height),
            image_path=image_path,
            opacity=opacity,
            scale_percent=scale_percent,
            position=position,
        ),
    )


def apply_text_overlay(
    input_pdf: Path,
    output_pdf: Path,
    text: str,
    page_spec: str = "",
    font_size: int = 24,
    rotation: float = 0.0,
    opacity: float = 1.0,
    position: str = "top-left",
) -> Path:
    if not text.strip():
        raise ConversionError("Overlay text cannot be empty.")

    reader = PdfReader(str(input_pdf))
    opacity = _clamp(float(opacity), 0.05, 1.0)
    target_pages = resolve_target_pages(page_spec, len(reader.pages), default="all")

    return _apply_overlay_to_pages(
        reader,
        output_pdf,
        target_pages,
        lambda _page_index, page: _build_text_watermark_page(
            float(page.mediabox.width),
            float(page.mediabox.height),
            text=text,
            font_size=font_size,
            rotation=rotation,
            opacity=opacity,
            position=position,
            color=(0.12, 0.12, 0.12),
        ),
    )


def apply_image_overlay(
    input_pdf: Path,
    output_pdf: Path,
    image_path: Path,
    page_spec: str = "",
    opacity: float = 1.0,
    scale_percent: int = 20,
    position: str = "top-right",
) -> Path:
    image_path = Path(image_path)
    if not image_path.exists():
        raise ConversionError(f"Overlay image was not found: {image_path}")

    reader = PdfReader(str(input_pdf))
    opacity = _clamp(float(opacity), 0.05, 1.0)
    scale_percent = max(int(scale_percent), 1)
    target_pages = resolve_target_pages(page_spec, len(reader.pages), default="all")

    return _apply_overlay_to_pages(
        reader,
        output_pdf,
        target_pages,
        lambda _page_index, page: _build_image_watermark_page(
            float(page.mediabox.width),
            float(page.mediabox.height),
            image_path=image_path,
            opacity=opacity,
            scale_percent=scale_percent,
            position=position,
        ),
    )


def _build_visible_signature_page(
    page_width: float,
    page_height: float,
    signer_text: str,
    signature_image: Path | None,
    opacity: float,
    scale_percent: int,
    position: str,
):
    buffer = BytesIO()
    overlay = canvas.Canvas(buffer, pagesize=(page_width, page_height))
    opacity = _clamp(opacity, 0.15, 1.0)

    box_width = min(page_width * 0.38, 260)
    box_height = min(page_height * 0.18, 110)
    x, y = _anchor_rect(page_width, page_height, box_width, box_height, position)

    overlay.saveState()
    if hasattr(overlay, "setFillAlpha"):
        try:
            overlay.setFillAlpha(opacity)
        except Exception:
            pass
    overlay.setFillColorRGB(1, 1, 1)
    overlay.setStrokeColorRGB(0.16, 0.36, 0.76)
    overlay.setLineWidth(1.1)
    overlay.roundRect(x, y, box_width, box_height, radius=10, fill=1, stroke=1)
    overlay.restoreState()

    content_x = x + 12
    content_y = y + box_height - 18
    text_x = content_x

    if signature_image:
        with Image.open(signature_image) as source_image:
            image = source_image.convert("RGBA")
            if opacity < 1.0:
                alpha = image.getchannel("A")
                alpha = alpha.point(lambda pixel: int(pixel * opacity))
                image.putalpha(alpha)
            image_reader = ImageReader(image)
            img_width, img_height = image.size
            available_width = box_width * 0.42
            available_height = box_height - 24
            scaled_width = min(available_width, max(48.0, page_width * (max(scale_percent, 1) / 100.0) * 0.35))
            scaled_height = scaled_width * (img_height / max(img_width, 1))
            if scaled_height > available_height:
                shrink = available_height / scaled_height
                scaled_width *= shrink
                scaled_height *= shrink
            img_x = x + 10
            img_y = y + (box_height - scaled_height) / 2
            overlay.drawImage(
                image_reader,
                img_x,
                img_y,
                width=scaled_width,
                height=scaled_height,
                mask="auto",
                preserveAspectRatio=True,
            )
            text_x = img_x + scaled_width + 12

    overlay.setFillColorRGB(0.12, 0.12, 0.12)
    overlay.setFont("Helvetica-Bold", 10)
    primary = signer_text.strip() or "Signed document"
    overlay.drawString(text_x, content_y, primary[:72])
    overlay.setFont("Helvetica", 8)
    overlay.setFillColorRGB(0.28, 0.28, 0.28)
    overlay.drawString(text_x, content_y - 16, f"Visible sign-off • {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    overlay.drawString(text_x, content_y - 30, "Generated by Gokul Omni Convert Lite")
    overlay.save()
    buffer.seek(0)
    return PdfReader(buffer).pages[0]


def apply_visible_signature(
    input_pdf: Path,
    output_pdf: Path,
    signer_text: str = "",
    signature_image: Path | None = None,
    page_spec: str = "",
    opacity: float = 0.95,
    scale_percent: int = 28,
    position: str = "bottom-right",
) -> Path:
    if signature_image is not None:
        signature_image = Path(signature_image)
        if not signature_image.exists():
            raise ConversionError(f"Signature image was not found: {signature_image}")

    if not signer_text.strip() and signature_image is None:
        raise ConversionError("Provide signature text, a signature image, or both.")

    reader = PdfReader(str(input_pdf))
    target_pages = resolve_target_pages(page_spec, len(reader.pages), default="last")

    return _apply_overlay_to_pages(
        reader,
        output_pdf,
        target_pages,
        lambda _page_index, page: _build_visible_signature_page(
            float(page.mediabox.width),
            float(page.mediabox.height),
            signer_text=signer_text,
            signature_image=signature_image,
            opacity=opacity,
            scale_percent=scale_percent,
            position=position,
        ),
    )




def _coerce_rect_value(raw: str, axis_length: float) -> float:
    value = str(raw).strip()
    if not value:
        raise ConversionError("Rectangle values cannot be blank.")
    if value.endswith("%"):
        return axis_length * (float(value[:-1]) / 100.0)
    numeric = float(value)
    if 0.0 <= numeric <= 1.0:
        return axis_length * numeric
    return numeric


def parse_rect_spec(rect_spec: str, page_rect: fitz.Rect) -> fitz.Rect:
    raw = str(rect_spec or "").strip()
    if not raw:
        raise ConversionError("Enter an area rectangle such as 36,72,420,160 or 10%,10%,90%,25%.")
    parts = [part for part in re.split(r"[\s,;]+", raw) if part]
    if len(parts) != 4:
        raise ConversionError("Area rectangle must contain exactly four values: x1, y1, x2, y2.")
    try:
        x1 = _coerce_rect_value(parts[0], page_rect.width)
        y1 = _coerce_rect_value(parts[1], page_rect.height)
        x2 = _coerce_rect_value(parts[2], page_rect.width)
        y2 = _coerce_rect_value(parts[3], page_rect.height)
    except ValueError as exc:
        raise ConversionError("Area rectangle values must be numbers or percentages.") from exc
    rect = fitz.Rect(x1, y1, x2, y2)
    rect = rect.normalize()
    page_box = fitz.Rect(page_rect)
    rect = fitz.Rect(
        max(page_box.x0, rect.x0),
        max(page_box.y0, rect.y0),
        min(page_box.x1, rect.x1),
        min(page_box.y1, rect.y1),
    )
    if rect.width <= 1 or rect.height <= 1:
        raise ConversionError("The selected rectangle is too small or falls outside the page bounds.")
    return rect


def redact_pdf_area(
    input_pdf: Path,
    output_pdf: Path,
    rect_spec: str,
    *,
    page_spec: str = "",
) -> tuple[Path, int]:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    with fitz.open(str(input_pdf)) as document:
        if document.needs_pass:
            raise ConversionError("This PDF is password protected. Unlock it first, then apply area redaction.")
        target_pages = resolve_target_pages(page_spec, len(document))
        redaction_count = 0
        for page_index in target_pages:
            page = document[page_index]
            redaction_rect = parse_rect_spec(rect_spec, page.rect)
            page.add_redact_annot(redaction_rect, fill=(0, 0, 0), cross_out=False)
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            redaction_count += 1
        document.save(str(output_pdf), garbage=4, deflate=True)
    return output_pdf, redaction_count


def edit_pdf_text_best_effort(
    input_pdf: Path,
    output_pdf: Path,
    *,
    search_text: str,
    replacement_text: str,
    page_spec: str = "",
) -> tuple[Path, int]:
    search_value = str(search_text or "").strip()
    if not search_value:
        raise ConversionError("Enter the text to find before running best-effort PDF text edit.")
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    with fitz.open(str(input_pdf)) as document:
        if document.needs_pass:
            raise ConversionError("This PDF is password protected. Unlock it first, then run text edit.")
        target_pages = resolve_target_pages(page_spec, len(document))
        match_count = 0
        for page_index in target_pages:
            page = document[page_index]
            matches = page.search_for(search_value)
            for rect in matches:
                font_size = max(8, min(18, rect.height * 0.75))
                page.add_redact_annot(
                    rect,
                    text=replacement_text,
                    fontname="helv",
                    fontsize=font_size,
                    align=fitz.TEXT_ALIGN_LEFT,
                    fill=(1, 1, 1),
                    text_color=(0, 0, 0),
                    cross_out=False,
                )
            if matches:
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
                match_count += len(matches)
        if match_count == 0:
            raise ConversionError("No extractable text matched the search term. Try text overlay for a non-destructive edit.")
        document.save(str(output_pdf), garbage=4, deflate=True)
    return output_pdf, match_count


def redact_pdf_text(
    input_pdf: Path,
    output_pdf: Path,
    search_text: str,
    page_spec: str = "",
) -> tuple[Path, int]:
    if not search_text.strip():
        raise ConversionError("Search text for redaction cannot be empty.")

    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)
    total_matches = 0

    with fitz.open(str(input_pdf)) as document:
        target_pages = resolve_target_pages(page_spec, document.page_count, default="all")
        for page_index in target_pages:
            page = document[page_index]
            matches = page.search_for(search_text)
            total_matches += len(matches)
            for rect in matches:
                page.add_redact_annot(rect, fill=(0, 0, 0))
            if matches:
                page.apply_redactions()
        if total_matches == 0:
            raise ConversionError(f"No matches were found for redaction text: {search_text}")
        document.save(str(output_pdf), garbage=4, deflate=True)

    return output_pdf, total_matches


def edit_pdf_metadata(
    input_pdf: Path,
    output_pdf: Path,
    *,
    title: str = "",
    author: str = "",
    subject: str = "",
    keywords: str = "",
    clear_existing: bool = False,
) -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    blank = {
        "title": "",
        "author": "",
        "subject": "",
        "keywords": "",
        "creator": "",
        "producer": "",
        "creationDate": "",
        "modDate": "",
        "trapped": "",
    }

    with fitz.open(str(input_pdf)) as document:
        metadata = dict(blank)
        if not clear_existing:
            for key in blank:
                metadata[key] = str(document.metadata.get(key, "") or "")
        metadata["title"] = title.strip() if clear_existing or title.strip() else metadata.get("title", "")
        metadata["author"] = author.strip() if clear_existing or author.strip() else metadata.get("author", "")
        metadata["subject"] = subject.strip() if clear_existing or subject.strip() else metadata.get("subject", "")
        metadata["keywords"] = keywords.strip() if clear_existing or keywords.strip() else metadata.get("keywords", "")
        if clear_existing and hasattr(document, "del_xml_metadata"):
            try:
                document.del_xml_metadata()
            except Exception:
                pass
        document.set_metadata(metadata)
        document.save(str(output_pdf), garbage=4, deflate=True)

    return output_pdf


def lock_pdf(input_pdf: Path, output_pdf: Path, user_password: str, owner_password: str = "") -> Path:
    if not user_password.strip():
        raise ConversionError("A password is required to lock the PDF.")

    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    with fitz.open(str(input_pdf)) as document:
        if document.needs_pass:
            raise ConversionError("This PDF is already password protected. Unlock it first before applying a new password.")
        permissions = int(
            fitz.PDF_PERM_ACCESSIBILITY
            | fitz.PDF_PERM_PRINT
            | fitz.PDF_PERM_COPY
            | fitz.PDF_PERM_ANNOTATE
            | fitz.PDF_PERM_FORM
            | fitz.PDF_PERM_PRINT_HQ
        )
        document.save(
            str(output_pdf),
            garbage=3,
            deflate=True,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=(owner_password or user_password)[:40],
            user_pw=user_password[:40],
            permissions=permissions,
        )
    return output_pdf


def unlock_pdf(input_pdf: Path, output_pdf: Path, password: str) -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    with fitz.open(str(input_pdf)) as document:
        if document.needs_pass:
            if not password:
                raise ConversionError("Enter the current PDF password to unlock this file.")
            if document.authenticate(password) <= 0:
                raise ConversionError("The supplied PDF password is incorrect.")
        document.save(
            str(output_pdf),
            garbage=3,
            deflate=True,
            encryption=fitz.PDF_ENCRYPT_NONE,
        )
    return output_pdf


def compress_pdf(input_pdf: Path, output_pdf: Path, profile: str = "balanced", password: str = "") -> Path:
    output_pdf = unique_path(output_pdf)
    ensure_directory(output_pdf.parent)

    profile_key = (profile or "balanced").strip().lower()
    profile_options = {
        "safe": dict(garbage=3, deflate=True),
        "balanced": dict(garbage=4, deflate=True, deflate_images=True, deflate_fonts=True),
        "strong": dict(garbage=4, clean=True, deflate=True, deflate_images=True, deflate_fonts=True),
    }
    options = profile_options.get(profile_key, profile_options["balanced"])

    with fitz.open(str(input_pdf)) as document:
        if document.needs_pass:
            if not password:
                raise ConversionError("This PDF is password protected. Enter the current password to compress it.")
            if document.authenticate(password) <= 0:
                raise ConversionError("The supplied PDF password is incorrect.")
        document.save(str(output_pdf), **options)
    return output_pdf


def process_batch(config: BatchConfig, log: LogFn | None = None, progress: ProgressFn | None = None) -> list[Path]:
    files = dedupe_paths(config.files)
    if not files:
        raise ConversionError("No input files selected.")

    ensure_directory(config.output_dir)
    total = len(files)
    outputs: list[Path] = []

    if config.mode == MODE_MERGE_PDFS:
        _log(log, f"Merging {total} PDF file(s)...")
        output_pdf = config.output_dir / f"{safe_name(config.merged_output_name or default_merged_name(config.mode))}.pdf"
        merged = merge_pdfs(files, output_pdf)
        _progress(progress, total, total)
        _log(log, f"Created: {merged}")
        return [merged]

    if config.mode == MODE_PDF_TO_IMAGES:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Converting PDF to images: {pdf_path.name}")
            target_dir = ensure_directory(config.output_dir / f"{safe_name(pdf_path.stem)}_images")
            pdf_to_images(pdf_path, target_dir, image_format=config.image_format, image_scale=config.image_scale)
            outputs.append(target_dir)
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_PDF_TO_DOCX:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Converting PDF to DOCX: {pdf_path.name}")
            output_docx = config.output_dir / f"{safe_name(pdf_path.stem)}.docx"
            outputs.append(pdf_to_docx(pdf_path, output_docx))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_PDF_TO_XLSX:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Converting PDF to XLSX: {pdf_path.name}")
            output_xlsx = config.output_dir / f"{safe_name(pdf_path.stem)}.xlsx"
            outputs.append(pdf_to_xlsx(pdf_path, output_xlsx))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_MD_TO_DOCX:
        for index, md_path in enumerate(files, start=1):
            _log(log, f"Converting Markdown to DOCX: {md_path.name}")
            output_docx = config.output_dir / f"{safe_name(md_path.stem)}.docx"
            outputs.append(markdown_to_docx(md_path, output_docx))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_MD_TO_HTML:
        for index, md_path in enumerate(files, start=1):
            _log(log, f"Converting Markdown to HTML: {md_path.name}")
            output_html = config.output_dir / f"{safe_name(md_path.stem)}.html"
            outputs.append(markdown_to_html(md_path, output_html))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_HTML_TO_DOCX:
        for index, html_path in enumerate(files, start=1):
            _log(log, f"Converting HTML to DOCX: {html_path.name}")
            output_docx = config.output_dir / f"{safe_name(html_path.stem)}.docx"
            outputs.append(html_to_docx(html_path, output_docx))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_HTML_TO_MD:
        for index, html_path in enumerate(files, start=1):
            _log(log, f"Converting HTML to Markdown: {html_path.name}")
            output_md = config.output_dir / f"{safe_name(html_path.stem)}.md"
            outputs.append(html_to_markdown(html_path, output_md))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_PDF_TO_HTML:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Converting PDF to HTML: {pdf_path.name}")
            output_html = config.output_dir / f"{safe_name(pdf_path.stem)}.html"
            outputs.append(pdf_to_html(pdf_path, output_html))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_PDF_TO_PPTX:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Converting PDF to PPTX: {pdf_path.name}")
            output_pptx = config.output_dir / f"{safe_name(pdf_path.stem)}.pptx"
            outputs.append(pdf_to_pptx(pdf_path, output_pptx, image_scale=config.image_scale))
            _progress(progress, index, total)
        return outputs

    if config.mode == MODE_PRESENTATIONS_TO_IMAGES:
        for index, presentation_path in enumerate(files, start=1):
            _log(log, f"Converting presentation to images: {presentation_path.name}")
            created = presentation_to_images(
                presentation_path,
                config.output_dir,
                image_format=config.image_format,
                image_scale=config.image_scale,
                engine_mode=config.engine_mode,
                soffice_path=config.soffice_path,
            )
            outputs.extend(created)
            for item in created:
                _log(log, f"Created: {item}")
            _progress(progress, index, total)
        return outputs

    if config.mode in PDF_OUTPUT_MODES:
        merge_to_one = bool(config.merge_to_one_pdf)

        if merge_to_one:
            output_pdf = config.output_dir / f"{safe_name(config.merged_output_name or default_merged_name(config.mode))}.pdf"
            _log(log, f"Building merged PDF from {total} file(s)...")

            if config.mode == MODE_IMAGES_TO_PDF and all(path.suffix.lower() in IMAGE_EXTS for path in files):
                merged = convert_images_to_single_pdf(files, output_pdf)
                _progress(progress, total, total)
                _log(log, f"Created: {merged}")
                return [merged]

            with tempfile.TemporaryDirectory() as temp_dir:
                temp_pdfs: list[Path] = []
                temp_root = Path(temp_dir)
                for index, input_path in enumerate(files, start=1):
                    temp_pdf = temp_root / f"{index:04d}_{safe_name(input_path.stem)}.pdf"
                    _log(log, f"Converting for merge: {input_path.name}")
                    _log(log, describe_input_route_for_mode(config.mode, input_path, engine_mode=config.engine_mode, soffice_path=config.soffice_path))
                    convert_input_to_pdf_for_mode(
                        config.mode,
                        input_path,
                        temp_pdf,
                        engine_mode=config.engine_mode,
                        soffice_path=config.soffice_path,
                    )
                    temp_pdfs.append(temp_pdf)
                    _progress(progress, index, total)
                merged = merge_pdfs(temp_pdfs, output_pdf)
                _log(log, f"Created: {merged}")
                return [merged]

        for index, input_path in enumerate(files, start=1):
            _log(log, f"Converting to PDF: {input_path.name}")
            _log(log, describe_input_route_for_mode(config.mode, input_path, engine_mode=config.engine_mode, soffice_path=config.soffice_path))
            output_pdf = config.output_dir / f"{safe_name(input_path.stem)}.pdf"
            outputs.append(
                convert_input_to_pdf_for_mode(
                    config.mode,
                    input_path,
                    output_pdf,
                    engine_mode=config.engine_mode,
                    soffice_path=config.soffice_path,
                )
            )
            _progress(progress, index, total)
        return outputs

    raise ConversionError(f"Unsupported mode: {config.mode}")


def process_pdf_tool(config: PdfToolConfig, log: LogFn | None = None, progress: ProgressFn | None = None) -> list[Path]:
    files = dedupe_paths(config.files)
    if not files:
        raise ConversionError("No PDF files selected for the PDF tool.")
    if any(path.suffix.lower() not in PDF_EXTS for path in files):
        raise ConversionError("PDF tools only accept PDF inputs.")

    ensure_directory(config.output_dir)
    total = len(files)
    outputs: list[Path] = []

    if config.tool == PDF_TOOL_MERGE:
        output_pdf = config.output_dir / f"{safe_name(config.output_name or 'merged_pdfs')}.pdf"
        _log(log, f"Merging {total} PDF file(s)...")
        merged = merge_pdfs(files, output_pdf)
        _progress(progress, total, total)
        _log(log, f"Created: {merged}")
        return [merged]

    if config.tool == PDF_TOOL_SPLIT_RANGES:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Splitting by ranges: {pdf_path.name}")
            created = split_pdf_by_ranges(pdf_path, config.page_spec, config.output_dir)
            outputs.extend(created)
            for path in created:
                _log(log, f"Created: {path}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_SPLIT_EVERY_N:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Splitting every {config.every_n_pages} page(s): {pdf_path.name}")
            created = split_pdf_every_n(pdf_path, config.every_n_pages, config.output_dir)
            outputs.extend(created)
            for path in created:
                _log(log, f"Created: {path}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_EXTRACT_PAGES:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Extracting pages from: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_extracted.pdf"
            outputs.append(extract_pdf_pages(pdf_path, config.page_spec, output_pdf))
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_REMOVE_PAGES:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Removing pages from: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_removed.pdf"
            outputs.append(remove_pdf_pages(pdf_path, config.page_spec, output_pdf))
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_REORDER_PAGES:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Reordering pages for: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_reordered.pdf"
            outputs.append(reorder_pdf_pages(pdf_path, config.page_spec, output_pdf))
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_WATERMARK_TEXT:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Adding text watermark to: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_watermarked.pdf"
            outputs.append(
                apply_text_watermark(
                    pdf_path,
                    output_pdf,
                    text=config.watermark_text,
                    font_size=config.watermark_font_size,
                    rotation=config.watermark_rotation,
                    opacity=config.watermark_opacity,
                    position=config.watermark_position,
                    page_spec=config.page_spec,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_WATERMARK_IMAGE:
        if not config.watermark_image:
            raise ConversionError("Please choose a watermark image file first.")
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Adding image watermark to: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_watermarked.pdf"
            outputs.append(
                apply_image_watermark(
                    pdf_path,
                    output_pdf,
                    image_path=config.watermark_image,
                    opacity=config.watermark_opacity,
                    scale_percent=config.watermark_image_scale_percent,
                    position=config.watermark_position,
                    page_spec=config.page_spec,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_TEXT_OVERLAY:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Adding text overlay to: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_edited_text.pdf"
            outputs.append(
                apply_text_overlay(
                    pdf_path,
                    output_pdf,
                    text=config.watermark_text,
                    page_spec=config.page_spec,
                    font_size=config.watermark_font_size,
                    rotation=config.watermark_rotation,
                    opacity=config.watermark_opacity,
                    position=config.watermark_position,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_IMAGE_OVERLAY:
        if not config.watermark_image:
            raise ConversionError("Please choose an image file for the PDF image overlay.")
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Adding image overlay to: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_edited_image.pdf"
            outputs.append(
                apply_image_overlay(
                    pdf_path,
                    output_pdf,
                    image_path=config.watermark_image,
                    page_spec=config.page_spec,
                    opacity=config.watermark_opacity,
                    scale_percent=config.watermark_image_scale_percent,
                    position=config.watermark_position,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_REDACT_TEXT:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Redacting searched text in: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_redacted.pdf"
            redacted_path, matches = redact_pdf_text(
                pdf_path,
                output_pdf,
                search_text=config.watermark_text,
                page_spec=config.page_spec,
            )
            outputs.append(redacted_path)
            _log(log, f"Redacted {matches} match(es) -> {redacted_path}")
            _progress(progress, index, total)
        return outputs


    if config.tool == PDF_TOOL_REDACT_AREA:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Redacting area in: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_redacted_area.pdf"
            redacted_path, matches = redact_pdf_area(
                pdf_path,
                output_pdf,
                rect_spec=config.redact_rect,
                page_spec=config.page_spec,
            )
            outputs.append(redacted_path)
            _log(log, f"Redacted {matches} page region(s) -> {redacted_path}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_EDIT_TEXT:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Running best-effort text edit in: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_edited_replace.pdf"
            edited_path, matches = edit_pdf_text_best_effort(
                pdf_path,
                output_pdf,
                search_text=config.watermark_text,
                replacement_text=config.replacement_text,
                page_spec=config.page_spec,
            )
            outputs.append(edited_path)
            _log(log, f"Replaced {matches} match(es) -> {edited_path}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_SIGN_VISIBLE:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Applying visible signature to: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_signed.pdf"
            outputs.append(
                apply_visible_signature(
                    pdf_path,
                    output_pdf,
                    signer_text=config.watermark_text,
                    signature_image=config.watermark_image,
                    page_spec=config.page_spec,
                    opacity=config.watermark_opacity,
                    scale_percent=config.watermark_image_scale_percent,
                    position=config.watermark_position,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_EDIT_METADATA:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Editing metadata for: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_metadata.pdf"
            outputs.append(
                edit_pdf_metadata(
                    pdf_path,
                    output_pdf,
                    title=config.metadata_title,
                    author=config.metadata_author,
                    subject=config.metadata_subject,
                    keywords=config.metadata_keywords,
                    clear_existing=config.metadata_clear_existing,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_LOCK:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Locking PDF with password: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_locked.pdf"
            outputs.append(
                lock_pdf(
                    pdf_path,
                    output_pdf,
                    user_password=config.pdf_password,
                    owner_password=config.pdf_owner_password,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_UNLOCK:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Unlocking PDF: {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_unlocked.pdf"
            outputs.append(unlock_pdf(pdf_path, output_pdf, password=config.pdf_password))
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    if config.tool == PDF_TOOL_COMPRESS:
        for index, pdf_path in enumerate(files, start=1):
            _log(log, f"Compressing PDF ({config.compression_profile}): {pdf_path.name}")
            output_pdf = config.output_dir / f"{safe_name(pdf_path.stem)}_compressed.pdf"
            outputs.append(
                compress_pdf(
                    pdf_path,
                    output_pdf,
                    profile=config.compression_profile,
                    password=config.pdf_password,
                )
            )
            _log(log, f"Created: {outputs[-1]}")
            _progress(progress, index, total)
        return outputs

    raise ConversionError(f"Unsupported PDF tool: {config.tool}")
