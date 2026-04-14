from __future__ import annotations

import csv
import re
from dataclasses import dataclass
from html import escape
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable, Sequence

from docx import Document
from docx.document import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as DocxTable, _Cell
from docx.text.paragraph import Paragraph as DocxParagraph
from docx.shared import Inches as DocxInches
from openpyxl import load_workbook
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, StyleSheet1, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import LongTable, PageBreak, Paragraph, Preformatted, SimpleDocTemplate, Spacer, Table, TableStyle


PAGE_WIDTH, PAGE_HEIGHT = A4
PAGE_LEFT_RIGHT_MARGIN = 14 * mm
PAGE_TOP_BOTTOM_MARGIN = 14 * mm
CONTENT_WIDTH = PAGE_WIDTH - (2 * PAGE_LEFT_RIGHT_MARGIN)


@dataclass(slots=True)
class HtmlBlock:
    type: str
    text: str = ""
    level: int = 0
    indent: int = 0
    bullet: str = ""
    rows: list[list[str]] | None = None


_HTML_BLOCK_TAGS = {"p", "div", "section", "article", "header", "footer"}
_INLINE_OPEN = {
    "strong": "<b>",
    "b": "<b>",
    "em": "<i>",
    "i": "<i>",
    "u": "<u>",
    "code": '<font face="Courier">',
}
_INLINE_CLOSE = {
    "strong": "</b>",
    "b": "</b>",
    "em": "</i>",
    "i": "</i>",
    "u": "</u>",
    "code": "</font>",
}


def _build_styles() -> StyleSheet1:
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="PureBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.5,
            leading=13.5,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="PureTitle",
            parent=styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=20,
            leading=24,
            spaceAfter=10,
            alignment=TA_CENTER,
            textColor=colors.HexColor("#0f172a"),
        )
    )
    styles.add(
        ParagraphStyle(
            name="PureSubtitle",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=9,
            leading=11,
            spaceAfter=12,
            alignment=TA_CENTER,
            textColor=colors.HexColor("#475569"),
        )
    )
    for level, size, after in [(1, 18, 10), (2, 15.5, 9), (3, 13.5, 8), (4, 12.5, 7), (5, 11.5, 6), (6, 10.8, 6)]:
        styles.add(
            ParagraphStyle(
                name=f"PureHeading{level}",
                parent=styles["Heading1"],
                fontName="Helvetica-Bold",
                fontSize=size,
                leading=size + 3,
                spaceBefore=10 if level <= 2 else 8,
                spaceAfter=after,
                textColor=colors.HexColor("#111827"),
            )
        )
    styles.add(
        ParagraphStyle(
            name="PureList",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.5,
            leading=13.5,
            leftIndent=14,
            firstLineIndent=0,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="PureQuote",
            parent=styles["BodyText"],
            fontName="Helvetica-Oblique",
            fontSize=10.2,
            leading=13,
            leftIndent=18,
            rightIndent=8,
            borderPadding=6,
            borderColor=colors.HexColor("#cbd5e1"),
            borderLeftWidth=3,
            backColor=colors.HexColor("#f8fafc"),
            textColor=colors.HexColor("#334155"),
            spaceBefore=4,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="PureCode",
            parent=styles["BodyText"],
            fontName="Courier",
            fontSize=9,
            leading=11,
            leftIndent=8,
            rightIndent=8,
            borderPadding=8,
            backColor=colors.HexColor("#111827"),
            textColor=colors.HexColor("#f8fafc"),
            spaceBefore=4,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="PureTableCell",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=8.5,
            leading=10.5,
            spaceAfter=0,
        )
    )
    styles.add(
        ParagraphStyle(
            name="PureTableHeader",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=8.6,
            leading=10.6,
            textColor=colors.white,
            spaceAfter=0,
        )
    )
    return styles


STYLES = _build_styles()


def _safe_markup(text: object) -> str:
    value = "" if text is None else str(text)
    return escape(value, quote=False).replace("\n", "<br/>")


def _simplify_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _coerce_rows(rows: Iterable[Sequence[object]]) -> list[list[str]]:
    result: list[list[str]] = []
    for row in rows:
        values = ["" if value is None else str(value) for value in row]
        if any(_simplify_text(value) for value in values):
            result.append(values)
    return result


def _column_widths_for_rows(rows: Sequence[Sequence[str]], available_width: float = CONTENT_WIDTH) -> list[float] | None:
    if not rows:
        return None
    column_count = max(len(row) for row in rows)
    if column_count < 1:
        return None
    weights = [1.0] * column_count
    for column_index in range(column_count):
        lengths: list[int] = []
        for row in rows:
            if column_index >= len(row):
                continue
            text = _simplify_text(row[column_index])
            lengths.append(min(max(len(text), 4), 28))
        if lengths:
            weights[column_index] = max(sum(lengths) / len(lengths), 6)
    total = sum(weights) or float(column_count)
    return [available_width * (weight / total) for weight in weights]


def _build_flow_table(rows: Sequence[Sequence[str]]) -> Table:
    materialized_rows = [list(row) for row in rows]
    widths = _column_widths_for_rows(materialized_rows)
    cell_style = STYLES["PureTableCell"]
    header_style = STYLES["PureTableHeader"]
    table_data: list[list[Paragraph]] = []
    for row_index, row in enumerate(materialized_rows):
        rendered_row: list[Paragraph] = []
        for cell in row:
            style = header_style if row_index == 0 else cell_style
            rendered_row.append(Paragraph(_safe_markup(cell) or "&nbsp;", style))
        table_data.append(rendered_row)
    table = LongTable(table_data, colWidths=widths, repeatRows=1 if len(table_data) > 1 else 0)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#334155")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#cbd5e1")),
                ("BOX", (0, 0), (-1, -1), 0.45, colors.HexColor("#94a3b8")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
                ("LEFTPADDING", (0, 0), (-1, -1), 5),
                ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return table


class _StructuredHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.blocks: list[HtmlBlock] = []
        self._current_text: list[str] = []
        self._current_kind = "paragraph"
        self._current_level = 0
        self._current_indent = 0
        self._current_bullet = ""
        self._list_stack: list[dict[str, object]] = []
        self._in_pre = False
        self._pre_parts: list[str] = []
        self._in_table = False
        self._table_rows: list[list[str]] = []
        self._current_row: list[str] = []
        self._current_cell: list[str] = []
        self._cell_has_markup = False
        self._saw_text = False
        self._ignored_depth = 0

    def _flush_text_block(self) -> None:
        raw_text = "".join(self._current_text).strip()
        if raw_text:
            self.blocks.append(
                HtmlBlock(
                    type=self._current_kind,
                    text=raw_text,
                    level=self._current_level,
                    indent=self._current_indent,
                    bullet=self._current_bullet,
                )
            )
        self._current_text = []
        self._current_kind = "paragraph"
        self._current_level = 0
        self._current_indent = 0
        self._current_bullet = ""

    def _flush_pre_block(self) -> None:
        text = "".join(self._pre_parts).rstrip("\n")
        if text:
            self.blocks.append(HtmlBlock(type="pre", text=text))
        self._pre_parts = []
        self._in_pre = False

    def handle_starttag(self, tag: str, attrs) -> None:  # type: ignore[override]
        tag = tag.lower()
        attrs_dict = dict(attrs or [])
        if tag in {"head", "style", "script", "title"}:
            self._ignored_depth += 1
            return
        if self._ignored_depth:
            return
        if tag == "br":
            if self._in_pre:
                self._pre_parts.append("\n")
            elif self._in_table and self._current_cell is not None:
                self._current_cell.append("<br/>")
            else:
                self._current_text.append("<br/>")
            return
        if tag in _HTML_BLOCK_TAGS:
            self._flush_text_block()
            self._current_kind = "paragraph"
            return
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._flush_text_block()
            self._current_kind = "heading"
            self._current_level = int(tag[1])
            return
        if tag == "blockquote":
            self._flush_text_block()
            self._current_kind = "quote"
            return
        if tag == "ul":
            self._flush_text_block()
            self._list_stack.append({"kind": "ul", "counter": 1})
            return
        if tag == "ol":
            self._flush_text_block()
            self._list_stack.append({"kind": "ol", "counter": 1})
            return
        if tag == "li":
            self._flush_text_block()
            level = len(self._list_stack)
            current = self._list_stack[-1] if self._list_stack else {"kind": "ul", "counter": 1}
            if current["kind"] == "ol":
                bullet = f"{current['counter']}."
                current["counter"] = int(current["counter"]) + 1
            else:
                bullet = "•"
            self._current_kind = "list_item"
            self._current_indent = max(level - 1, 0)
            self._current_bullet = bullet
            return
        if tag == "pre":
            self._flush_text_block()
            self._in_pre = True
            self._pre_parts = []
            return
        if tag == "table":
            self._flush_text_block()
            if self._in_pre:
                self._flush_pre_block()
            self._in_table = True
            self._table_rows = []
            return
        if tag == "tr" and self._in_table:
            self._current_row = []
            return
        if tag in {"td", "th"} and self._in_table:
            self._current_cell = []
            self._cell_has_markup = False
            return
        if tag == "a":
            href = _safe_markup(attrs_dict.get("href", "")).strip()
            if href:
                self._current_text.append('<font color="#2563eb"><u>')
                self._saw_text = True
            return
        if tag in _INLINE_OPEN:
            target = self._current_cell if self._in_table and self._current_cell is not None else self._current_text
            target.append(_INLINE_OPEN[tag])
            self._saw_text = True
            return

    def handle_endtag(self, tag: str) -> None:  # type: ignore[override]
        tag = tag.lower()
        if tag in {"head", "style", "script", "title"}:
            if self._ignored_depth:
                self._ignored_depth -= 1
            return
        if self._ignored_depth:
            return
        if tag in _HTML_BLOCK_TAGS | {"li", "blockquote"}:
            self._flush_text_block()
            return
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._flush_text_block()
            return
        if tag in {"ul", "ol"}:
            self._flush_text_block()
            if self._list_stack:
                self._list_stack.pop()
            return
        if tag == "pre":
            self._flush_pre_block()
            return
        if tag in {"td", "th"} and self._in_table:
            cell_text = "".join(self._current_cell).strip() if self._current_cell else ""
            self._current_row.append(cell_text)
            self._current_cell = []
            return
        if tag == "tr" and self._in_table:
            if any(_simplify_text(value) for value in self._current_row):
                self._table_rows.append(self._current_row)
            self._current_row = []
            return
        if tag == "table":
            if self._table_rows:
                self.blocks.append(HtmlBlock(type="table", rows=[list(row) for row in self._table_rows]))
            self._table_rows = []
            self._in_table = False
            return
        if tag == "a":
            self._current_text.append("</u></font>")
            return
        if tag in _INLINE_CLOSE:
            target = self._current_cell if self._in_table and self._current_cell is not None else self._current_text
            target.append(_INLINE_CLOSE[tag])
            return

    def handle_data(self, data: str) -> None:  # type: ignore[override]
        if self._ignored_depth:
            return
        if not data:
            return
        if self._in_pre:
            self._pre_parts.append(data)
            return
        target = self._current_cell if self._in_table and self._current_cell is not None else self._current_text
        cleaned = data
        if not self._in_table:
            cleaned = re.sub(r"[\t\f\v ]+", " ", data)
        text = escape(cleaned, quote=False)
        if text:
            target.append(text)
            self._saw_text = True

    def get_blocks(self) -> list[HtmlBlock]:
        if self._in_pre:
            self._flush_pre_block()
        else:
            self._flush_text_block()
        if self._in_table and self._table_rows:
            self.blocks.append(HtmlBlock(type="table", rows=[list(row) for row in self._table_rows]))
            self._in_table = False
        return [block for block in self.blocks if block.text or block.rows]


def _html_to_blocks(html_content: str) -> list[HtmlBlock]:
    parser = _StructuredHTMLParser()
    parser.feed(html_content)
    parser.close()
    blocks = parser.get_blocks()
    return blocks


def _table_blocks_to_story(rows: Sequence[Sequence[str]]) -> list[object]:
    if not rows:
        return []
    column_count = max(len(row) for row in rows)
    if column_count <= 8:
        return [_build_flow_table(rows), Spacer(1, 8)]

    story: list[object] = []
    for start in range(0, column_count, 8):
        end = min(start + 8, column_count)
        chunked = []
        for row in rows:
            chunked.append(list(row[start:end]))
        story.append(Paragraph(f"Columns {start + 1}-{end}", STYLES["PureSubtitle"]))
        story.append(_build_flow_table(chunked))
        story.append(Spacer(1, 8))
    return story


def render_html_to_pdf_from_string(html_content: str, output_pdf: Path, title: str = "") -> Path:
    output_pdf = Path(output_pdf)
    output_pdf.parent.mkdir(parents=True, exist_ok=True)

    blocks = _html_to_blocks(html_content)
    document_title = _simplify_text(title)
    story: list[object] = []
    if document_title:
        story.append(Paragraph(_safe_markup(document_title), STYLES["PureTitle"]))
        story.append(Spacer(1, 4))

    for block in blocks:
        if block.type == "heading":
            style = STYLES[f"PureHeading{min(max(block.level, 1), 6)}"]
            story.append(Paragraph(block.text, style))
            continue
        if block.type == "paragraph":
            story.append(Paragraph(block.text, STYLES["PureBody"]))
            continue
        if block.type == "quote":
            story.append(Paragraph(block.text, STYLES["PureQuote"]))
            continue
        if block.type == "list_item":
            list_style = ParagraphStyle(
                name=f"PureListDepth{block.indent}",
                parent=STYLES["PureList"],
                leftIndent=16 + (block.indent * 14),
                firstLineIndent=0,
            )
            story.append(Paragraph(block.text, list_style, bulletText=block.bullet or "•"))
            continue
        if block.type == "pre":
            story.append(Preformatted(block.text, STYLES["PureCode"]))
            continue
        if block.type == "table" and block.rows:
            story.extend(_table_blocks_to_story(block.rows))
            continue

    if not story:
        placeholder = document_title or "Converted document"
        story = [Paragraph(_safe_markup(placeholder), STYLES["PureTitle"]), Paragraph("No extractable HTML content was found.", STYLES["PureBody"])]

    doc = SimpleDocTemplate(
        str(output_pdf),
        pagesize=A4,
        leftMargin=PAGE_LEFT_RIGHT_MARGIN,
        rightMargin=PAGE_LEFT_RIGHT_MARGIN,
        topMargin=PAGE_TOP_BOTTOM_MARGIN,
        bottomMargin=PAGE_TOP_BOTTOM_MARGIN,
        title=document_title or output_pdf.stem,
    )
    doc.build(story)
    return output_pdf


def _iter_block_items(parent: DocxDocument | _Cell):
    if isinstance(parent, DocxDocument):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise TypeError(f"Unsupported parent type: {type(parent)!r}")
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield DocxParagraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)


def _paragraph_markup(paragraph: DocxParagraph) -> str:
    if paragraph.runs:
        parts: list[str] = []
        for run in paragraph.runs:
            text = _safe_markup(run.text)
            if not text:
                continue
            if run.bold:
                text = f"<b>{text}</b>"
            if run.italic:
                text = f"<i>{text}</i>"
            if run.underline:
                text = f"<u>{text}</u>"
            parts.append(text)
        combined = "".join(parts).strip()
        if combined:
            return combined
    return _safe_markup(paragraph.text)


def _table_rows_from_docx(table: DocxTable) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in table.rows:
        values: list[str] = []
        for cell in row.cells:
            text_lines: list[str] = []
            for block in _iter_block_items(cell):
                if isinstance(block, DocxParagraph):
                    value = _simplify_text(block.text)
                    if value:
                        text_lines.append(value)
                elif isinstance(block, DocxTable):
                    text_lines.append("[Nested table omitted]")
            values.append("\n".join(text_lines).strip())
        rows.append(values)
    return _coerce_rows(rows)


def render_docx_to_pdf(input_docx: Path, output_pdf: Path) -> Path:
    output_pdf = Path(output_pdf)
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    document = Document(str(input_docx))
    story: list[object] = [Paragraph(_safe_markup(input_docx.stem), STYLES["PureTitle"]), Spacer(1, 6)]
    extracted_any = False

    for block in _iter_block_items(document):
        if isinstance(block, DocxParagraph):
            raw_text = _simplify_text(block.text)
            if not raw_text:
                continue
            style_name = (getattr(getattr(block, "style", None), "name", "") or "").lower()
            markup = _paragraph_markup(block)
            if style_name.startswith("heading"):
                digits = re.findall(r"\d+", style_name)
                level = int(digits[0]) if digits else 1
                story.append(Paragraph(markup, STYLES[f"PureHeading{min(max(level, 1), 6)}"]))
            elif "quote" in style_name:
                story.append(Paragraph(markup, STYLES["PureQuote"]))
            elif "list" in style_name or raw_text.startswith(("•", "-", "*")):
                bullet_text = "•"
                text_body = markup
                match = re.match(r"^(?:•|-|\*)\s+(.*)$", raw_text)
                if match:
                    text_body = _safe_markup(match.group(1))
                list_style = ParagraphStyle(name="PureDocxList", parent=STYLES["PureList"], leftIndent=18)
                story.append(Paragraph(text_body, list_style, bulletText=bullet_text))
            else:
                story.append(Paragraph(markup, STYLES["PureBody"]))
            extracted_any = True
        elif isinstance(block, DocxTable):
            rows = _table_rows_from_docx(block)
            if rows:
                story.extend(_table_blocks_to_story(rows))
                extracted_any = True

    if document.inline_shapes:
        story.append(Spacer(1, 6))
        story.append(
            Paragraph(
                f"Embedded images detected: {len(document.inline_shapes)}. Pure Python export keeps the text and table structure, while image placement may differ from Word.",
                STYLES["PureSubtitle"],
            )
        )

    if not extracted_any:
        story.append(Paragraph("No extractable text or tables were found in this DOCX file.", STYLES["PureBody"]))

    doc = SimpleDocTemplate(
        str(output_pdf),
        pagesize=A4,
        leftMargin=PAGE_LEFT_RIGHT_MARGIN,
        rightMargin=PAGE_LEFT_RIGHT_MARGIN,
        topMargin=PAGE_TOP_BOTTOM_MARGIN,
        bottomMargin=PAGE_TOP_BOTTOM_MARGIN,
        title=input_docx.stem,
    )
    doc.build(story)
    return output_pdf


def _load_spreadsheet_rows(input_path: Path) -> list[tuple[str, list[list[str]]]]:
    ext = input_path.suffix.lower()
    sheets: list[tuple[str, list[list[str]]]] = []
    if ext == ".xlsx":
        workbook = load_workbook(filename=str(input_path), data_only=True, read_only=True)
        try:
            for sheet in workbook.worksheets:
                rows = _coerce_rows(sheet.iter_rows(values_only=True))
                sheets.append((sheet.title, rows))
        finally:
            workbook.close()
        return sheets
    if ext == ".xls":
        import xlrd  # type: ignore

        book = xlrd.open_workbook(str(input_path), on_demand=True)
        try:
            for sheet in book.sheets():
                rows: list[list[str]] = []
                for row_index in range(sheet.nrows):
                    values = [sheet.cell_value(row_index, col_index) for col_index in range(sheet.ncols)]
                    rows.append(["" if value is None else str(value) for value in values])
                sheets.append((sheet.name, _coerce_rows(rows)))
        finally:
            book.release_resources()
        return sheets
    if ext in {".csv", ".tsv"}:
        delimiter = "\t" if ext == ".tsv" else ","
        with input_path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            reader = csv.reader(handle, delimiter=delimiter)
            sheets.append((input_path.stem, _coerce_rows(reader)))
        return sheets
    raise ValueError(f"Unsupported spreadsheet input: {input_path.name}")


def render_spreadsheet_to_pdf(input_sheet: Path, output_pdf: Path) -> Path:
    output_pdf = Path(output_pdf)
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    sheets = _load_spreadsheet_rows(input_sheet)
    story: list[object] = [Paragraph(_safe_markup(input_sheet.stem), STYLES["PureTitle"]), Spacer(1, 6)]

    for index, (sheet_name, rows) in enumerate(sheets, start=1):
        if index > 1:
            story.append(PageBreak())
        story.append(Paragraph(_safe_markup(sheet_name or f"Sheet {index}"), STYLES["PureHeading2"]))
        if rows:
            story.extend(_table_blocks_to_story(rows))
        else:
            story.append(Paragraph("This sheet does not contain visible cell values.", STYLES["PureBody"]))

    if not sheets:
        story.append(Paragraph("No readable sheets were found in this workbook.", STYLES["PureBody"]))

    doc = SimpleDocTemplate(
        str(output_pdf),
        pagesize=A4,
        leftMargin=PAGE_LEFT_RIGHT_MARGIN,
        rightMargin=PAGE_LEFT_RIGHT_MARGIN,
        topMargin=PAGE_TOP_BOTTOM_MARGIN,
        bottomMargin=PAGE_TOP_BOTTOM_MARGIN,
        title=input_sheet.stem,
    )
    doc.build(story)
    return output_pdf


def render_presentation_to_pdf(input_pptx: Path, output_pdf: Path) -> Path:
    from pptx import Presentation

    output_pdf = Path(output_pdf)
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(str(input_pptx))
    story: list[object] = [Paragraph(_safe_markup(input_pptx.stem), STYLES["PureTitle"]), Spacer(1, 6)]

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > 1:
            story.append(PageBreak())
        story.append(Paragraph(f"Slide {slide_index}", STYLES["PureHeading2"]))
        text_blocks: list[str] = []
        image_count = 0
        for shape in slide.shapes:
            text_value = getattr(shape, "text", "")
            if text_value and _simplify_text(text_value):
                text_blocks.append(text_value.strip())
            if getattr(shape, "shape_type", None) == 13:  # picture
                image_count += 1
        if text_blocks:
            for text_value in text_blocks:
                story.append(Paragraph(_safe_markup(text_value), STYLES["PureBody"]))
        else:
            story.append(Paragraph("No extractable text was found on this slide.", STYLES["PureBody"]))
        if image_count:
            story.append(Paragraph(f"Images detected on this slide: {image_count}. Pure Python export preserves slide text and notes the image count.", STYLES["PureSubtitle"]))

    doc = SimpleDocTemplate(
        str(output_pdf),
        pagesize=A4,
        leftMargin=PAGE_LEFT_RIGHT_MARGIN,
        rightMargin=PAGE_LEFT_RIGHT_MARGIN,
        topMargin=PAGE_TOP_BOTTOM_MARGIN,
        bottomMargin=PAGE_TOP_BOTTOM_MARGIN,
        title=input_pptx.stem,
    )
    doc.build(story)
    return output_pdf


def render_html_to_docx_from_string(html_content: str, output_docx: Path, title: str = "") -> Path:
    output_docx = Path(output_docx)
    output_docx.parent.mkdir(parents=True, exist_ok=True)
    blocks = _html_to_blocks(html_content)
    document = Document()
    if title:
        document.add_heading(title, level=0)
    for block in blocks:
        if block.type == "heading":
            document.add_heading(re.sub(r"<[^>]+>", "", block.text), level=min(max(block.level, 1), 4))
        elif block.type == "paragraph":
            document.add_paragraph(re.sub(r"<[^>]+>", "", block.text))
        elif block.type == "quote":
            paragraph = document.add_paragraph(re.sub(r"<[^>]+>", "", block.text))
            try:
                paragraph.style = "Intense Quote"
            except Exception:
                pass
        elif block.type == "list_item":
            paragraph = document.add_paragraph(style="List Bullet" if not re.match(r"^\d+\.$", block.bullet) else "List Number")
            paragraph.add_run(re.sub(r"<[^>]+>", "", block.text))
        elif block.type == "pre":
            paragraph = document.add_paragraph()
            run = paragraph.add_run(block.text)
            run.font.name = "Courier New"
        elif block.type == "table" and block.rows:
            row_count = len(block.rows)
            column_count = max(len(row) for row in block.rows)
            table = document.add_table(rows=row_count, cols=column_count)
            for row_index, row in enumerate(block.rows):
                for column_index in range(column_count):
                    value = row[column_index] if column_index < len(row) else ""
                    table.cell(row_index, column_index).text = value
    document.save(str(output_docx))
    return output_docx
