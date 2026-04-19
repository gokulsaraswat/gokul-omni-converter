from __future__ import annotations

import html
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import fitz
from PIL import Image, ImageDraw, ImageFont, ImageOps

try:
    from docx import Document
except Exception:  # pragma: no cover - optional at runtime
    Document = None  # type: ignore

try:
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - optional at runtime
    load_workbook = None  # type: ignore

try:
    import xlrd
except Exception:  # pragma: no cover - optional at runtime
    xlrd = None  # type: ignore

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional at runtime
    Presentation = None  # type: ignore


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
TEXT_EXTENSIONS = {".txt", ".log", ".json", ".xml", ".yaml", ".yml", ".ini", ".cfg", ".py", ".csv", ".tsv"}
MARKDOWN_EXTENSIONS = {".md", ".markdown"}
HTML_EXTENSIONS = {".html", ".htm"}
DOC_EXTENSIONS = {".docx", ".doc", ".odt", ".rtf"}
SHEET_EXTENSIONS = {".xlsx", ".xls", ".xlsm", ".csv", ".tsv", ".ods"}
PRESENTATION_EXTENSIONS = {".pptx", ".ppt", ".odp"}


@dataclass(slots=True)
class PreviewResult:
    image: Image.Image
    summary: str
    title: str
    kind: str
    page_count: int = 1
    current_page: int = 0


def _font(size: int, *, bold: bool = False) -> ImageFont.ImageFont:
    candidates = []
    if bold:
        candidates.extend(["DejaVuSans-Bold.ttf", "Arial Bold.ttf", "arialbd.ttf"])
    candidates.extend(["DejaVuSans.ttf", "Arial.ttf", "arial.ttf"])
    for name in candidates:
        try:
            return ImageFont.truetype(name, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = text.replace("\r", "").split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        probe = f"{current} {word}".strip()
        if draw.textlength(probe, font=font) <= max_width:
            current = probe
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def _sanitize_excerpt(text: str) -> str:
    cleaned = text.replace("\u00a0", " ").replace("\r", "\n")
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _read_text_file(path: Path, max_chars: int = 6000) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    if len(raw) > max_chars:
        raw = raw[:max_chars] + "\n..."
    return _sanitize_excerpt(raw)


def _html_to_text(html_text: str) -> str:
    text = re.sub(r"<script[\s\S]*?</script>", " ", html_text, flags=re.IGNORECASE)
    text = re.sub(r"<style[\s\S]*?</style>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s+", "\n", text)
    text = re.sub(r"\s{2,}", " ", text)
    return _sanitize_excerpt(text)


def _docx_excerpt(path: Path, max_lines: int = 18) -> list[str]:
    if Document is None:
        return ["python-docx is unavailable in this environment."]
    doc = Document(str(path))
    lines: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
        if len(lines) >= max_lines:
            break
    if not lines:
        lines.append("No paragraph text was found in this document.")
    return lines


def _sheet_excerpt(path: Path, max_rows: int = 8, max_cols: int = 6) -> list[str]:
    ext = path.suffix.lower()
    lines: list[str] = []
    if ext == ".xls":
        if xlrd is None:
            return ["xlrd is unavailable in this environment."]
        workbook = xlrd.open_workbook(str(path))
        sheet = workbook.sheet_by_index(0)
        lines.append(f"Sheet: {sheet.name}")
        for row_index in range(min(sheet.nrows, max_rows)):
            row_values = []
            for col_index in range(min(sheet.ncols, max_cols)):
                value = sheet.cell_value(row_index, col_index)
                row_values.append(str(value))
            lines.append(" | ".join(row_values))
        return lines
    if ext in {".csv", ".tsv"}:
        delimiter = "\t" if ext == ".tsv" else ","
        text = _read_text_file(path, max_chars=3000)
        rows = [row for row in text.splitlines() if row.strip()]
        for row in rows[:max_rows]:
            parts = [part.strip() for part in row.split(delimiter)[:max_cols]]
            lines.append(" | ".join(parts))
        return lines or ["No rows found in this file."]
    if load_workbook is None:
        return ["openpyxl is unavailable in this environment."]
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    sheet = workbook[workbook.sheetnames[0]]
    lines.append(f"Sheet: {sheet.title}")
    for row in sheet.iter_rows(min_row=1, max_row=max_rows, values_only=True):
        parts = ["" if value is None else str(value) for value in row[:max_cols]]
        if any(part.strip() for part in parts):
            lines.append(" | ".join(parts))
    if len(lines) == 1:
        lines.append("No populated cells were found in the first rows.")
    return lines


def _presentation_excerpt(path: Path, max_slides: int = 6) -> list[str]:
    if path.suffix.lower() != ".pptx" or Presentation is None:
        return ["Preview summary is available for PPTX files. Other presentation formats show metadata only."]
    presentation = Presentation(str(path))
    lines = [f"Slides: {len(presentation.slides)}"]
    for slide_index, slide in enumerate(list(presentation.slides)[:max_slides], start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None:
                for paragraph in text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)
        joined = " • ".join(texts[:3]) or "(no text found on this slide)"
        lines.append(f"Slide {slide_index}: {joined}")
    return lines


def _generic_summary(path: Path) -> list[str]:
    stat = path.stat()
    size_kb = stat.st_size / 1024 if stat.st_size else 0
    return [
        f"File name: {path.name}",
        f"Type: {path.suffix.lower() or 'unknown'}",
        f"Size: {size_kb:.1f} KB",
        "A richer visual preview is not available for this format, but the app can still process it when supported.",
    ]


def _render_card(*, title: str, subtitle: str, lines: Iterable[str], footer: str = "", size: tuple[int, int] = (1200, 1500)) -> Image.Image:
    width, height = size
    image = Image.new("RGB", size, "#ffffff")
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((24, 24, width - 24, height - 24), radius=32, fill="#f8fafc", outline="#dbe4ee", width=2)
    draw.rounded_rectangle((48, 48, width - 48, 190), radius=24, fill="#0f172a")
    title_font = _font(40, bold=True)
    subtitle_font = _font(24)
    body_font = _font(24)
    small_font = _font(20)

    draw.text((80, 76), title[:120], fill="#f8fafc", font=title_font)
    draw.text((80, 132), subtitle[:180], fill="#cbd5e1", font=subtitle_font)

    y = 240
    max_width = width - 140
    for raw_line in lines:
        line = (raw_line or "").strip()
        if not line:
            y += 12
            continue
        wrapped = _wrap_text(draw, line, body_font, max_width)
        for piece in wrapped[:6]:
            draw.text((82, y), piece, fill="#0f172a", font=body_font)
            y += 34
            if y > height - 160:
                break
        if y > height - 160:
            break
    if footer:
        draw.line((72, height - 120, width - 72, height - 120), fill="#d7e1eb", width=2)
        footer_lines = _wrap_text(draw, footer, small_font, width - 140)
        footer_text = "\n".join(footer_lines[:3])
        draw.multiline_text((82, height - 100), footer_text, fill="#475569", font=small_font, spacing=6)
    return image


def _image_preview(path: Path) -> PreviewResult:
    with Image.open(path) as img:
        preview = ImageOps.exif_transpose(img).convert("RGB")
        result = preview.copy()
        summary = f"Image preview\nFile: {path.name}\nPixels: {result.width} x {result.height}\nFormat: {img.format or path.suffix.upper()}"
        return PreviewResult(image=result, summary=summary, title=path.name, kind="image")


def _pdf_preview(path: Path, *, page: int = 0, zoom: float = 1.0) -> PreviewResult:
    document = fitz.open(str(path))
    if document.page_count == 0:
        placeholder = _render_card(title=path.name, subtitle="Empty PDF", lines=["This PDF does not contain any pages."])
        return PreviewResult(image=placeholder, summary="Empty PDF.", title=path.name, kind="pdf", page_count=0, current_page=0)
    clamped_page = max(0, min(page, document.page_count - 1))
    page_obj = document.load_page(clamped_page)
    matrix = fitz.Matrix(max(0.5, min(zoom, 3.0)), max(0.5, min(zoom, 3.0)))
    pixmap = page_obj.get_pixmap(matrix=matrix, alpha=False)
    image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
    meta = document.metadata or {}
    summary_lines = [
        f"PDF preview",
        f"File: {path.name}",
        f"Page: {clamped_page + 1} / {document.page_count}",
    ]
    if meta.get("title"):
        summary_lines.append(f"Title: {meta.get('title')}")
    if meta.get("author"):
        summary_lines.append(f"Author: {meta.get('author')}")
    summary = "\n".join(summary_lines)
    return PreviewResult(image=image, summary=summary, title=path.name, kind="pdf", page_count=document.page_count, current_page=clamped_page)


def _text_based_preview(path: Path, kind: str, lines: list[str], subtitle: str) -> PreviewResult:
    card = _render_card(title=path.name, subtitle=subtitle, lines=lines[:24], footer=f"Path: {path}")
    summary = "\n".join([subtitle, f"File: {path.name}", f"Lines shown: {min(len(lines), 24)}"])
    return PreviewResult(image=card, summary=summary, title=path.name, kind=kind)


def _doc_preview(path: Path) -> PreviewResult:
    if path.suffix.lower() == ".docx":
        lines = _docx_excerpt(path)
        return _text_based_preview(path, "document", lines, "DOCX preview summary")
    return _text_based_preview(path, "document", _generic_summary(path), "Office document metadata")


def _sheet_preview(path: Path) -> PreviewResult:
    lines = _sheet_excerpt(path)
    return _text_based_preview(path, "sheet", lines, "Spreadsheet preview summary")


def _presentation_preview(path: Path) -> PreviewResult:
    lines = _presentation_excerpt(path)
    return _text_based_preview(path, "presentation", lines, "Presentation preview summary")


def _text_preview(path: Path) -> PreviewResult:
    ext = path.suffix.lower()
    raw = _read_text_file(path)
    text = _html_to_text(raw) if ext in HTML_EXTENSIONS else raw
    lines = [line for line in text.splitlines() if line.strip()][:28]
    subtitle = "HTML preview summary" if ext in HTML_EXTENSIONS else ("Markdown preview summary" if ext in MARKDOWN_EXTENSIONS else "Text preview summary")
    return _text_based_preview(path, "text", lines or ["(file is empty)"], subtitle)


def preview_supported(path: str | Path) -> bool:
    try:
        source = Path(path)
    except Exception:
        return False
    return source.suffix.lower() in (
        IMAGE_EXTENSIONS
        | TEXT_EXTENSIONS
        | MARKDOWN_EXTENSIONS
        | HTML_EXTENSIONS
        | DOC_EXTENSIONS
        | SHEET_EXTENSIONS
        | PRESENTATION_EXTENSIONS
        | {".pdf"}
    )


def render_preview(path: str | Path, *, page: int = 0, zoom: float = 1.0) -> PreviewResult:
    source = Path(path).expanduser()
    if not source.exists():
        placeholder = _render_card(title=source.name or "Missing file", subtitle="File unavailable", lines=[str(source), "This path does not exist anymore."])
        return PreviewResult(image=placeholder, summary=f"The file is missing:\n{source}", title=source.name or "Missing", kind="missing")
    ext = source.suffix.lower()
    if ext == ".pdf":
        return _pdf_preview(source, page=page, zoom=zoom)
    if ext in IMAGE_EXTENSIONS:
        return _image_preview(source)
    if ext in TEXT_EXTENSIONS or ext in MARKDOWN_EXTENSIONS or ext in HTML_EXTENSIONS:
        return _text_preview(source)
    if ext in DOC_EXTENSIONS:
        return _doc_preview(source)
    if ext in SHEET_EXTENSIONS:
        return _sheet_preview(source)
    if ext in PRESENTATION_EXTENSIONS:
        return _presentation_preview(source)
    return _text_based_preview(source, "generic", _generic_summary(source), "Generic file summary")


def preview_summary(path: str | Path) -> str:
    result = render_preview(path)
    return result.summary
